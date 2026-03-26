"""
Enhanced PPTX Toolset v4
------------------------
Key improvements over v3:
  - Gemini-powered keyword generation: analyzes full slide content to produce
    highly specific, photogenic search terms (not generic filler)
  - Multi-source image scraping: DuckDuckGo image search fetches N candidates
  - Gemini Vision selection: evaluates each candidate image against slide context
    and picks the most semantically relevant one (never just "first result")
  - No .venv required: uses system Python + stdlib urllib / requests
  - Fallback chain: DuckDuckGo -> Unsplash source -> Picsum

Image pipeline per slide:
  slide content (title + body + bullets)
       |  Gemini Flash text
       v
  5 specific search queries
       |  DuckDuckGo image scrape (3 URLs per query -> up to 15 candidates)
       v
  download top 5 candidates
       |  Gemini Flash Vision (send images as base64)
       v
  pick most contextually relevant image
"""

from __future__ import annotations

import base64
import io
import json
import logging
import os
import re
import urllib.parse
import urllib.request
from typing import Any

import matplotlib
matplotlib.use("Agg")
import matplotlib.pyplot as plt
import numpy as np
from PIL import Image

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

logger = logging.getLogger(__name__)

# ---------------------------------------------------------------------------
# Gemini API helpers
# ---------------------------------------------------------------------------

GEMINI_API_KEY = os.getenv("GEMINI_API_KEY", "")
_GEMINI_TEXT_URL = (
    "https://generativelanguage.googleapis.com/v1beta/models/"
    "gemini-2.0-flash:generateContent"
)
_IMG_TIMEOUT = 10
_HEADERS = {
    "User-Agent": (
        "Mozilla/5.0 (Windows NT 10.0; Win64; x64) "
        "AppleWebKit/537.36 (KHTML, like Gecko) "
        "Chrome/124.0 Safari/537.36"
    ),
    "Accept": "image/webp,image/apng,image/*,*/*;q=0.8",
}


def _gemini_post(payload: dict) -> dict | None:
    """Call Gemini API with a JSON payload. Returns parsed response or None."""
    if not GEMINI_API_KEY:
        return None
    url = f"{_GEMINI_TEXT_URL}?key={GEMINI_API_KEY}"
    data = json.dumps(payload).encode()
    req = urllib.request.Request(
        url,
        data=data,
        headers={"Content-Type": "application/json"},
        method="POST",
    )
    try:
        with urllib.request.urlopen(req, timeout=20) as resp:
            return json.loads(resp.read().decode())
    except Exception as exc:
        logger.warning(f"Gemini API error: {exc}")
        return None


def _gemini_generate_keywords(slide_data: dict) -> list[str]:
    """
    Ask Gemini to produce 5 highly specific, photogenic image-search queries
    tailored to the slide's actual content.
    Returns a list of query strings, falling back to the slide title if failed.
    """
    title   = slide_data.get("title", "")
    body    = slide_data.get("body", "")
    bullets = " | ".join(slide_data.get("bullets", [])[:3])
    quote   = slide_data.get("quote", "")
    subtitle = slide_data.get("subtitle", "")
    context = f"Slide title: {title}\nBody: {body}\nSubtitle: {subtitle}\nQuote: {quote}\nKey points: {bullets}"

    prompt = (
        "You are an expert visual researcher for presentations. "
        "Given this presentation slide content, generate exactly 5 HIGHLY SPECIFIC "
        "image search queries that would retrieve the most relevant, high-quality "
        "stock photos for this slide.\n\n"
        "Rules:\n"
        "- Be EXTREMELY specific: include subject, setting, action, mood, context\n"
        "- Think like a photo editor: what real-world image would illustrate this?\n"
        "- Avoid vague or abstract terms; think photogenic, concrete, real-world scenes\n"
        "- Each query must be completely distinct (different subject/angle/setting)\n"
        "- Each query must be 3-6 words\n"
        "- Queries must be directly relevant to the slide TOPIC, not generic business photos\n"
        "- Return ONLY a JSON array of 5 strings, no other text\n\n"
        f"Slide content:\n{context}\n\n"
        'Output format example: ["surgeon performing robotic surgery", "hospital operating room equipment", ...]'
    )

    resp = _gemini_post({
        "contents": [{"parts": [{"text": prompt}]}],
        "generationConfig": {"temperature": 0.2, "maxOutputTokens": 300},
    })

    if resp:
        try:
            raw = resp["candidates"][0]["content"]["parts"][0]["text"].strip()
            raw = re.sub(r"^```[a-z]*\n?", "", raw)
            raw = re.sub(r"\n?```$", "", raw)
            queries = json.loads(raw)
            if isinstance(queries, list) and queries:
                logger.info(f"Gemini keywords for '{title}': {queries}")
                return [str(q) for q in queries[:5]]
        except Exception as exc:
            logger.warning(f"Failed to parse Gemini keywords: {exc}")

    # Fallback: derive from title
    words = (title + " " + body).split()[:6]
    fallbacks = slide_data.get("fallback_keywords", [])
    base = [" ".join(words[:4]), title, " ".join(words[2:6])]
    return (base + fallbacks)[:5] if base[0] else ["professional business office"]


def _gemini_select_best_image(
    candidates: list[tuple[str, bytes]],
    slide_data: dict,
) -> bytes | None:
    """
    Send up to 5 candidate images to Gemini Vision.
    Ask it to pick the most contextually relevant one.
    Returns the bytes of the winner, or the first candidate if selection fails.
    """
    if not candidates:
        return None
    if len(candidates) == 1:
        return candidates[0][1]

    title   = slide_data.get("title", "")
    body    = slide_data.get("body", "")
    bullets = " | ".join(slide_data.get("bullets", [])[:4])
    quote   = slide_data.get("quote", "")
    context = f"Title: {title}. {body}. {quote}. Key points: {bullets}"

    parts: list[dict] = [
        {
            "text": (
                f"You are a presentation designer selecting the best stock photo.\n"
                f"SLIDE CONTEXT: {context}\n\n"
                f"I'm showing you {len(candidates)} candidate images (numbered 1-{len(candidates)}).\n"
                "You MUST pick the image that is most relevant and visually appropriate "
                "for this exact slide topic. Consider:\n"
                "1. Does the image subject directly relate to the slide topic?\n"
                "2. Is it professional quality and composition?\n"
                "3. Does it add visual meaning to the slide content?\n"
                "4. Avoid generic office/people shots when the topic is more specific.\n\n"
                "Reply with ONLY a single integer (1, 2, 3, etc.) — the number of the best image. "
                "No explanation."
            )
        }
    ]

    for i, (query, img_bytes) in enumerate(candidates, start=1):
        try:
            img = Image.open(io.BytesIO(img_bytes)).convert("RGB")
            img.thumbnail((640, 480))
            buf = io.BytesIO()
            img.save(buf, format="JPEG", quality=80)
            b64 = base64.b64encode(buf.getvalue()).decode()
            parts.append({"text": f"Image {i} (search used: {query}):"})
            parts.append({
                "inline_data": {
                    "mime_type": "image/jpeg",
                    "data": b64,
                }
            })
        except Exception as exc:
            logger.warning(f"Could not encode candidate {i}: {exc}")

    resp = _gemini_post({
        "contents": [{"parts": parts}],
        "generationConfig": {"temperature": 0.0, "maxOutputTokens": 10},
    })

    if resp:
        try:
            answer = resp["candidates"][0]["content"]["parts"][0]["text"].strip()
            match = re.search(r"\d+", answer)
            if match:
                idx = int(match.group()) - 1
                if 0 <= idx < len(candidates):
                    logger.info(
                        f"Gemini Vision picked image {idx + 1}/{len(candidates)} "
                        f"(query: '{candidates[idx][0]}') for slide '{slide_data.get('title', '')}'"
                    )
                    return candidates[idx][1]
        except Exception as exc:
            logger.warning(f"Gemini selection parse error: {exc}")

    logger.info("Gemini Vision unavailable — using first candidate")
    return candidates[0][1]


# ---------------------------------------------------------------------------
# Image scraping helpers
# ---------------------------------------------------------------------------

def _ddg_image_urls(query: str, max_results: int = 4) -> list[str]:
    """
    Use DuckDuckGo image search (unofficial) to get direct image URLs.
    Returns up to max_results image URLs.
    """
    urls: list[str] = []
    try:
        search_url = (
            "https://duckduckgo.com/?q="
            + urllib.parse.quote(query)
            + "&iax=images&ia=images"
        )
        req = urllib.request.Request(
            search_url,
            headers={"User-Agent": _HEADERS["User-Agent"]}
        )
        with urllib.request.urlopen(req, timeout=_IMG_TIMEOUT) as r:
            html = r.read().decode("utf-8", errors="replace")

        # Extract vqd token
        vqd_match = re.search(r'vqd=(["\'])([^"\']+)\1', html)
        if not vqd_match:
            vqd_match = re.search(r"vqd=([\d\-]+)", html)
        if not vqd_match:
            return []

        vqd = vqd_match.group(2) if len(vqd_match.groups()) >= 2 else vqd_match.group(1)

        # Fetch image results JSON
        api_url = (
            "https://duckduckgo.com/i.js?q="
            + urllib.parse.quote(query)
            + "&vqd=" + urllib.parse.quote(vqd)
            + "&p=1&s=0&u=bing&f=,,,&l=en-us"
        )
        req2 = urllib.request.Request(
            api_url,
            headers={
                "User-Agent": _HEADERS["User-Agent"],
                "Referer": "https://duckduckgo.com/",
                "Accept": "application/json",
            },
        )
        with urllib.request.urlopen(req2, timeout=_IMG_TIMEOUT) as r2:
            data = json.loads(r2.read().decode())

        for result in data.get("results", [])[:max_results]:
            img_url = result.get("image", "")
            if img_url and img_url.startswith("http"):
                urls.append(img_url)

    except Exception as exc:
        logger.debug(f"DuckDuckGo search failed for '{query}': {exc}")

    return urls


def _unsplash_url(query: str, width: int = 1920, height: int = 1080) -> str:
    kw = urllib.parse.quote(query, safe="")
    return f"https://source.unsplash.com/{width}x{height}/?{kw}"


def _download_image(url: str) -> bytes | None:
    """Download image bytes from a URL, with basic quality checks."""
    try:
        req = urllib.request.Request(url, headers=_HEADERS)
        with urllib.request.urlopen(req, timeout=_IMG_TIMEOUT) as resp:
            data = resp.read()
        if len(data) < 8_000:
            return None
        # Verify it's a valid image
        Image.open(io.BytesIO(data)).verify()
        return data
    except Exception:
        return None


def _fetch_best_image(
    slide_data: dict,
    width: int = 1920,
    height: int = 1080,
    extra_keywords: list[str] | None = None,
) -> bytes | None:
    """
    Full Gemini-powered image pipeline:
    1. Ask Gemini to generate 5 specific search queries for this slide
    2. Scrape DuckDuckGo for candidate image URLs (up to 4 per query)
    3. Download one successful image per query (up to 5 total candidates)
    4. Ask Gemini Vision to pick the most relevant candidate
    5. Return winning image bytes

    Falls back to Unsplash source API, then Picsum if all else fails.
    """
    # Step 1: Generate keywords via Gemini
    queries = _gemini_generate_keywords(slide_data)
    if extra_keywords:
        queries = (queries + extra_keywords)[:7]

    candidates: list[tuple[str, bytes]] = []

    # Step 2 & 3: Scrape and download candidates
    for query in queries:
        if len(candidates) >= 5:
            break
        # Try DuckDuckGo first
        img_urls = _ddg_image_urls(query, max_results=4)
        for img_url in img_urls:
            img_bytes = _download_image(img_url)
            if img_bytes:
                candidates.append((query, img_bytes))
                break  # one good image per query

        # If DDG yielded nothing for this query, try Unsplash source
        if not any(q == query for q, _ in candidates):
            unsplash = _download_image(_unsplash_url(query, width, height))
            if unsplash:
                candidates.append((query, unsplash))

    # If still empty, try Unsplash source with first 3 queries
    if not candidates:
        logger.warning("DuckDuckGo returned nothing — falling back to Unsplash source")
        for q in queries[:3]:
            img = _download_image(_unsplash_url(q, width, height))
            if img:
                candidates.append((q, img))
            if len(candidates) >= 3:
                break

    # Absolute last resort: Picsum
    if not candidates:
        title = slide_data.get("title", "")
        seed = abs(hash(title[:20])) % 5000
        img = _download_image(f"https://picsum.photos/seed/{seed}/{width}/{height}")
        if img:
            candidates.append(("picsum_fallback", img))

    if not candidates:
        return None

    # Step 4: Gemini Vision selects the best
    return _gemini_select_best_image(candidates, slide_data)


# Backwards-compatible wrapper for calls that pass a keyword string
def _fetch_stock_photo(
    keyword: str,
    width: int = 1920,
    height: int = 1080,
    fallback_keywords: list | None = None,
    _slide_data: dict | None = None,
) -> bytes | None:
    if _slide_data is None:
        _slide_data = {
            "title": keyword,
            "body": "",
            "bullets": fallback_keywords or [],
        }
    return _fetch_best_image(_slide_data, width, height, fallback_keywords)


def _bytes_stream(data: bytes) -> io.BytesIO:
    buf = io.BytesIO(data)
    buf.seek(0)
    return buf


def _styled_fallback_stream(
    width_px: int, height_px: int, accent: RGBColor, keyword: str = ""
) -> io.BytesIO:
    import colorsys
    r, g, b = accent.red / 255, accent.green / 255, accent.blue / 255
    h, s, v = colorsys.rgb_to_hsv(r, g, b)
    light_rgb = colorsys.hsv_to_rgb(h, max(0.10, s * 0.35), min(1.0, v * 1.25))
    fig_w, fig_h = width_px / 150, height_px / 150
    fig, ax = plt.subplots(figsize=(fig_w, fig_h), facecolor="white")
    ax.set_axis_off()
    grad = np.linspace(0, 1, 256).reshape(1, -1)
    grad = np.vstack([grad] * 64)
    ca, cb = [r, g, b], list(light_rgb)
    rgba = np.zeros((64, 256, 4))
    for ch in range(3):
        rgba[:, :, ch] = ca[ch] + (cb[ch] - ca[ch]) * grad
    rgba[:, :, 3] = 1.0
    ax.imshow(rgba, aspect="auto", extent=[0, 1, 0, 1], transform=ax.transAxes)
    if keyword:
        ax.text(0.5, 0.5, keyword.title(), transform=ax.transAxes,
                ha="center", va="center", fontsize=max(12, fig_w * 2),
                color="white", fontweight="bold", alpha=0.55, wrap=True)
    plt.tight_layout(pad=0)
    buf = io.BytesIO()
    fig.savefig(buf, format="PNG", dpi=150, bbox_inches="tight")
    plt.close(fig)
    buf.seek(0)
    return buf


# ---------------------------------------------------------------------------
# Themes - 12 professional palettes
# ---------------------------------------------------------------------------

THEMES: dict[str, dict] = {
    "blue": {
        "accent": RGBColor(0x1A, 0x4F, 0x8A),
        "mid":    RGBColor(0x26, 0x7A, 0xB5),
        "light":  RGBColor(0xE3, 0xEF, 0xF8),
        "bg":     RGBColor(0xFA, 0xFB, 0xFF),
        "white":  RGBColor(0xFF, 0xFF, 0xFF),
        "dark":   RGBColor(0x1A, 0x1A, 0x2E),
        "muted":  RGBColor(0x5C, 0x6B, 0x80),
        "mpl":    "#1A4F8A",
    },
    "green": {
        "accent": RGBColor(0x0D, 0x5C, 0x2E),
        "mid":    RGBColor(0x1D, 0x8A, 0x4A),
        "light":  RGBColor(0xD4, 0xED, 0xDA),
        "bg":     RGBColor(0xF6, 0xFB, 0xF7),
        "white":  RGBColor(0xFF, 0xFF, 0xFF),
        "dark":   RGBColor(0x0A, 0x1E, 0x12),
        "muted":  RGBColor(0x4A, 0x6A, 0x54),
        "mpl":    "#0D5C2E",
    },
    "dark": {
        "accent": RGBColor(0x0F, 0x1A, 0x35),
        "mid":    RGBColor(0x5B, 0x2D, 0x8E),
        "light":  RGBColor(0x2A, 0x2A, 0x45),
        "bg":     RGBColor(0x14, 0x14, 0x28),
        "white":  RGBColor(0xF0, 0xF0, 0xF8),
        "dark":   RGBColor(0xE8, 0xE8, 0xF5),
        "muted":  RGBColor(0x9A, 0x9A, 0xBB),
        "mpl":    "#5B2D8E",
    },
    "red": {
        "accent": RGBColor(0x8B, 0x14, 0x14),
        "mid":    RGBColor(0xC0, 0x2B, 0x2B),
        "light":  RGBColor(0xF9, 0xE0, 0xE0),
        "bg":     RGBColor(0xFF, 0xFA, 0xFA),
        "white":  RGBColor(0xFF, 0xFF, 0xFF),
        "dark":   RGBColor(0x1A, 0x08, 0x08),
        "muted":  RGBColor(0x7A, 0x55, 0x55),
        "mpl":    "#8B1414",
    },
    "purple": {
        "accent": RGBColor(0x4A, 0x0E, 0x7E),
        "mid":    RGBColor(0x7B, 0x2D, 0xBF),
        "light":  RGBColor(0xEA, 0xD9, 0xF7),
        "bg":     RGBColor(0xFB, 0xF7, 0xFF),
        "white":  RGBColor(0xFF, 0xFF, 0xFF),
        "dark":   RGBColor(0x18, 0x06, 0x2A),
        "muted":  RGBColor(0x6A, 0x55, 0x80),
        "mpl":    "#4A0E7E",
    },
    "orange": {
        "accent": RGBColor(0xC0, 0x4A, 0x08),
        "mid":    RGBColor(0xE8, 0x7A, 0x20),
        "light":  RGBColor(0xFD, 0xEA, 0xD5),
        "bg":     RGBColor(0xFF, 0xFB, 0xF7),
        "white":  RGBColor(0xFF, 0xFF, 0xFF),
        "dark":   RGBColor(0x1E, 0x0E, 0x00),
        "muted":  RGBColor(0x7A, 0x5A, 0x40),
        "mpl":    "#C04A08",
    },
    "teal": {
        "accent": RGBColor(0x02, 0x72, 0x80),
        "mid":    RGBColor(0x0A, 0xA8, 0xB8),
        "light":  RGBColor(0xCC, 0xF0, 0xF4),
        "bg":     RGBColor(0xF4, 0xFD, 0xFF),
        "white":  RGBColor(0xFF, 0xFF, 0xFF),
        "dark":   RGBColor(0x01, 0x20, 0x24),
        "muted":  RGBColor(0x40, 0x70, 0x76),
        "mpl":    "#027280",
    },
    "midnight": {
        "accent": RGBColor(0x0B, 0x24, 0x4F),
        "mid":    RGBColor(0xC9, 0xA0, 0x2A),
        "light":  RGBColor(0xE8, 0xEE, 0xF8),
        "bg":     RGBColor(0xF5, 0xF7, 0xFF),
        "white":  RGBColor(0xFF, 0xFF, 0xFF),
        "dark":   RGBColor(0x08, 0x14, 0x2E),
        "muted":  RGBColor(0x6A, 0x72, 0x90),
        "mpl":    "#0B244F",
    },
    "slate": {
        "accent": RGBColor(0x2E, 0x3A, 0x50),
        "mid":    RGBColor(0x48, 0xA9, 0xDC),
        "light":  RGBColor(0xE8, 0xEC, 0xF2),
        "bg":     RGBColor(0xF8, 0xF9, 0xFB),
        "white":  RGBColor(0xFF, 0xFF, 0xFF),
        "dark":   RGBColor(0x14, 0x1C, 0x2A),
        "muted":  RGBColor(0x7A, 0x85, 0x99),
        "mpl":    "#2E3A50",
    },
    "rose": {
        "accent": RGBColor(0x9C, 0x27, 0x50),
        "mid":    RGBColor(0xE0, 0x5C, 0x88),
        "light":  RGBColor(0xF9, 0xE0, 0xEB),
        "bg":     RGBColor(0xFF, 0xF7, 0xFA),
        "white":  RGBColor(0xFF, 0xFF, 0xFF),
        "dark":   RGBColor(0x28, 0x08, 0x16),
        "muted":  RGBColor(0x88, 0x5A, 0x6A),
        "mpl":    "#9C2750",
    },
    "forest": {
        "accent": RGBColor(0x1B, 0x3A, 0x24),
        "mid":    RGBColor(0xB8, 0x7A, 0x1E),
        "light":  RGBColor(0xD8, 0xE8, 0xD8),
        "bg":     RGBColor(0xF5, 0xFA, 0xF5),
        "white":  RGBColor(0xFF, 0xFF, 0xFF),
        "dark":   RGBColor(0x0A, 0x14, 0x0C),
        "muted":  RGBColor(0x5A, 0x6E, 0x5C),
        "mpl":    "#1B3A24",
    },
    "indigo": {
        "accent": RGBColor(0x31, 0x35, 0x9E),
        "mid":    RGBColor(0x60, 0x65, 0xDE),
        "light":  RGBColor(0xE2, 0xE3, 0xF8),
        "bg":     RGBColor(0xF7, 0xF7, 0xFF),
        "white":  RGBColor(0xFF, 0xFF, 0xFF),
        "dark":   RGBColor(0x0E, 0x10, 0x2E),
        "muted":  RGBColor(0x6A, 0x6C, 0xA8),
        "mpl":    "#31359E",
    },
}


# ---------------------------------------------------------------------------
# Chart builders
# ---------------------------------------------------------------------------

def _palette(accent_hex: str, n: int) -> list:
    import colorsys
    r = int(accent_hex[1:3], 16) / 255
    g = int(accent_hex[3:5], 16) / 255
    b = int(accent_hex[5:7], 16) / 255
    h, s, v = colorsys.rgb_to_hsv(r, g, b)
    out = []
    for i in range(n):
        ri, gi, bi = colorsys.hsv_to_rgb(
            (h + i * 0.08) % 1.0, max(0.25, s - i * 0.05), max(0.45, v - i * 0.04))
        out.append("#%02x%02x%02x" % (int(ri*255), int(gi*255), int(bi*255)))
    return out


def _chart_style(ax, title):
    ax.set_title(title, fontsize=14, fontweight="bold", color="#222", pad=14)
    ax.tick_params(axis="x", labelsize=10, colors="#555")
    ax.tick_params(axis="y", labelsize=9, colors="#888")
    ax.spines[["top", "right"]].set_visible(False)
    ax.spines[["left", "bottom"]].set_color("#DDD")
    ax.yaxis.grid(True, color="#F0F0F0", linewidth=0.8)
    ax.set_axisbelow(True)
    ax.set_facecolor("#FAFAFA")


def _chart_bar(labels, values, title, accent_hex, figsize=(10, 5)) -> io.BytesIO:
    colors = _palette(accent_hex, len(labels))
    fig, ax = plt.subplots(figsize=figsize, facecolor="white")
    bars = ax.bar(labels, values, color=colors, edgecolor="white", linewidth=1.0, width=0.6)
    _chart_style(ax, title)
    max_v = max(values) if values else 1
    for bar, val in zip(bars, values):
        label = f"{val:,}" if isinstance(val, int) else str(val)
        ax.text(bar.get_x() + bar.get_width() / 2, bar.get_height() + max_v * 0.012,
                label, ha="center", va="bottom", fontsize=9, color="#333", fontweight="bold")
    plt.tight_layout()
    buf = io.BytesIO()
    fig.savefig(buf, format="PNG", dpi=150)
    plt.close(fig)
    buf.seek(0)
    return buf


def _chart_line(labels, values, title, accent_hex, figsize=(10, 5)) -> io.BytesIO:
    fig, ax = plt.subplots(figsize=figsize, facecolor="white")
    ax.plot(labels, values, color=accent_hex, linewidth=2.5, marker="o",
            markersize=8, markerfacecolor="white",
            markeredgewidth=2.5, markeredgecolor=accent_hex)
    ax.fill_between(range(len(labels)), values, alpha=0.10, color=accent_hex)
    ax.set_xticks(range(len(labels)))
    ax.set_xticklabels(labels, fontsize=10, color="#555")
    _chart_style(ax, title)
    plt.tight_layout()
    buf = io.BytesIO()
    fig.savefig(buf, format="PNG", dpi=150)
    plt.close(fig)
    buf.seek(0)
    return buf


def _chart_pie(labels, values, title, accent_hex, figsize=(7, 4.5)) -> io.BytesIO:
    colors = _palette(accent_hex, len(labels))
    fig, ax = plt.subplots(figsize=figsize, facecolor="white")
    wedges, _, autotexts = ax.pie(
        values, labels=None, colors=colors,
        autopct="%1.1f%%", startangle=140,
        wedgeprops=dict(width=0.6, edgecolor="white", linewidth=2),
        pctdistance=0.77,
    )
    for at in autotexts:
        at.set_fontsize(9)
        at.set_color("white")
        at.set_fontweight("bold")
    ax.legend(wedges, labels, loc="center left",
              bbox_to_anchor=(1, 0, 0.5, 1), fontsize=9, frameon=False)
    ax.set_title(title, fontsize=14, fontweight="bold", color="#222", pad=14)
    plt.tight_layout()
    buf = io.BytesIO()
    fig.savefig(buf, format="PNG", dpi=150, bbox_inches="tight")
    plt.close(fig)
    buf.seek(0)
    return buf


def _chart_hbar(labels, values, title, accent_hex, figsize=(9, 5)) -> io.BytesIO:
    colors = _palette(accent_hex, len(labels))
    fig, ax = plt.subplots(figsize=figsize, facecolor="white")
    bars = ax.barh(labels[::-1], values[::-1], color=colors[::-1],
                   edgecolor="white", linewidth=0.8, height=0.55)
    ax.set_title(title, fontsize=14, fontweight="bold", color="#222", pad=12)
    ax.tick_params(axis="y", labelsize=11, colors="#444")
    ax.tick_params(axis="x", labelsize=9, colors="#888")
    ax.spines[["top", "right", "left"]].set_visible(False)
    ax.spines["bottom"].set_color("#DDD")
    ax.xaxis.grid(True, color="#F0F0F0", linewidth=0.8)
    ax.set_axisbelow(True)
    ax.set_facecolor("#FAFAFA")
    max_v = max(values) if values else 1
    for bar, val in zip(bars, values[::-1]):
        label = f"{val:,}" if isinstance(val, int) else str(val)
        ax.text(val + max_v * 0.01, bar.get_y() + bar.get_height() / 2,
                label, va="center", fontsize=9, color="#333", fontweight="bold")
    plt.tight_layout()
    buf = io.BytesIO()
    fig.savefig(buf, format="PNG", dpi=150)
    plt.close(fig)
    buf.seek(0)
    return buf


# ---------------------------------------------------------------------------
# Low-level drawing helpers
# ---------------------------------------------------------------------------

def _set_bg(slide, color: RGBColor):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def _rect(slide, left, top, width, height, color: RGBColor):
    shp = slide.shapes.add_shape(1, left, top, width, height)
    shp.fill.solid()
    shp.fill.fore_color.rgb = color
    shp.line.fill.background()
    return shp


def _textbox(slide, left, top, width, height, text, size, bold, color: RGBColor,
             align=PP_ALIGN.LEFT, wrap=True, italic=False):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = wrap
    tf.text = text
    for para in tf.paragraphs:
        para.alignment = align
        for run in para.runs:
            run.font.size = Pt(size)
            run.font.bold = bold
            run.font.italic = italic
            run.font.color.rgb = color
    return tb


def _header_band(slide, W, t, height=None):
    if height is None:
        height = Inches(1.25)
    _rect(slide, 0, 0, W, height, t["accent"])
    _rect(slide, 0, height - Inches(0.055), W, Inches(0.055), t["mid"])


def _footer_rule(slide, W, H, t):
    _rect(slide, 0, H - Inches(0.075), W, Inches(0.075), t["accent"])


# ---------------------------------------------------------------------------
# Slide builders - image slides pass full slide_data to Gemini pipeline
# ---------------------------------------------------------------------------

def _build_title_slide(prs, data, t):
    W, H = prs.slide_width, prs.slide_height
    sl = prs.slides.add_slide(prs.slide_layouts[6])
    keyword = data.get("keyword", "")
    used_image = False
    if keyword:
        img_bytes = _fetch_best_image(data, 1920, 1080, data.get("fallback_keywords"))
        if img_bytes:
            sl.shapes.add_picture(_bytes_stream(img_bytes), 0, 0, W, H)
            _rect(sl, 0, 0, int(W * 0.55), H, RGBColor(0x08, 0x08, 0x10))
            _rect(sl, int(W * 0.45), 0, int(W * 0.55), H, RGBColor(0x08, 0x08, 0x14))
            _rect(sl, 0, 0, Inches(0.35), H, t["mid"])
            used_image = True
    if not used_image:
        _set_bg(sl, t["accent"])
        _rect(sl, 0, 0, Inches(0.4), H, t["mid"])
        _rect(sl, 0, H - Inches(1.9), W, Inches(1.9), t["bg"])
    text_col = t["white"]
    sub_col = t["light"]
    _textbox(sl, Inches(0.85), Inches(1.1), W - Inches(1.7), Inches(2.8),
             data.get("title", "Presentation"), 42, True, text_col)
    sub = data.get("subtitle", "")
    if sub:
        _textbox(sl, Inches(0.85), Inches(4.1), W - Inches(1.7), Inches(1.1),
                 sub, 21, False, sub_col)
    presenter = data.get("presenter", "")
    if presenter:
        _textbox(sl, Inches(0.85), H - Inches(1.05), W - Inches(1.7), Inches(0.65),
                 presenter, 14, False,
                 RGBColor(0xBB, 0xBB, 0xCC) if used_image else t["muted"],
                 italic=True)


def _build_content_slide(prs, data, t):
    W, H = prs.slide_width, prs.slide_height
    sl = prs.slides.add_slide(prs.slide_layouts[6])
    _set_bg(sl, t["bg"])
    hdr_h = Inches(1.25)
    _header_band(sl, W, t, hdr_h)
    _textbox(sl, Inches(0.55), Inches(0.18), W - Inches(1.1), hdr_h - Inches(0.22),
             data.get("title", "Slide"), 26, True, t["white"])
    y = hdr_h + Inches(0.28)
    body = data.get("body", "")
    if body:
        _textbox(sl, Inches(0.6), y, W - Inches(1.2), Inches(1.25), body, 15, False, t["dark"])
        y += Inches(1.3)
    bullets = data.get("bullets", [])
    if bullets:
        tb = sl.shapes.add_textbox(Inches(0.6), y, W - Inches(1.2), H - y - Inches(0.55))
        tf = tb.text_frame
        tf.word_wrap = True
        for i, b in enumerate(bullets):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            p.text = f"\u25b8  {b}"
            p.space_before = Pt(8)
            for run in p.runs:
                run.font.size = Pt(17)
                run.font.color.rgb = t["dark"]
    notes = data.get("notes", "")
    if notes:
        sl.notes_slide.notes_text_frame.text = notes
    _footer_rule(sl, W, H, t)


def _build_two_column_slide(prs, data, t):
    W, H = prs.slide_width, prs.slide_height
    sl = prs.slides.add_slide(prs.slide_layouts[6])
    _set_bg(sl, t["bg"])
    hdr_h = Inches(1.25)
    _header_band(sl, W, t, hdr_h)
    _textbox(sl, Inches(0.55), Inches(0.18), W - Inches(1.1), hdr_h - Inches(0.22),
             data.get("title", "Comparison"), 26, True, t["white"])
    mid = W // 2
    col_top = hdr_h + Inches(0.22)
    col_h = H - col_top - Inches(0.55)
    col_w = mid - Inches(0.85)
    _rect(sl, mid - Inches(0.025), col_top, Inches(0.05), col_h, t["light"])

    def _col(lx, ctitle, cbody, cbullets):
        _rect(sl, lx, col_top, Inches(0.07), Inches(0.5), t["mid"])
        _textbox(sl, lx + Inches(0.15), col_top + Inches(0.04),
                 col_w - Inches(0.12), Inches(0.48), ctitle, 18, True, t["accent"])
        cy = col_top + Inches(0.6)
        if cbody:
            _textbox(sl, lx + Inches(0.1), cy, col_w - Inches(0.1), Inches(1.0),
                     cbody, 13, False, t["muted"])
            cy += Inches(1.05)
        if cbullets:
            tb = sl.shapes.add_textbox(lx + Inches(0.1), cy,
                                       col_w - Inches(0.1), col_h - (cy - col_top))
            tf = tb.text_frame
            tf.word_wrap = True
            for i, b in enumerate(cbullets):
                p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
                p.text = f"\u25b8  {b}"
                p.space_before = Pt(7)
                for run in p.runs:
                    run.font.size = Pt(15)
                    run.font.color.rgb = t["dark"]

    _col(Inches(0.5), data.get("left_title", "Option A"),
         data.get("left_body", ""), data.get("left_bullets", []))
    _col(mid + Inches(0.35), data.get("right_title", "Option B"),
         data.get("right_body", ""), data.get("right_bullets", []))
    _footer_rule(sl, W, H, t)


def _build_closing_slide(prs, data, t):
    W, H = prs.slide_width, prs.slide_height
    sl = prs.slides.add_slide(prs.slide_layouts[6])
    _set_bg(sl, t["accent"])
    _rect(sl, 0, 0, Inches(0.45), H, t["mid"])
    _rect(sl, W - Inches(3.8), H - Inches(3.8), Inches(3.8), Inches(3.8), t["mid"])
    _rect(sl, 0, H - Inches(2.1), W, Inches(2.1), t["mid"])
    _textbox(sl, Inches(1.3), Inches(1.5), W - Inches(2.2), Inches(2.9),
             data.get("title", "Thank You"), 50, True, t["white"], align=PP_ALIGN.CENTER)
    sub = data.get("subtitle", "")
    if sub:
        _textbox(sl, Inches(1.3), Inches(4.6), W - Inches(2.2), Inches(0.95),
                 sub, 20, False, t["light"], align=PP_ALIGN.CENTER)
    contact = data.get("contact", "")
    if contact:
        _textbox(sl, Inches(1.3), H - Inches(1.65), W - Inches(2.2), Inches(0.65),
                 contact, 13, False, t["light"], align=PP_ALIGN.CENTER, italic=True)


def _build_image_slide(prs, data, t):
    """Full-bleed photo + title overlay. Uses Gemini pipeline."""
    W, H = prs.slide_width, prs.slide_height
    sl = prs.slides.add_slide(prs.slide_layouts[6])
    img_bytes = _fetch_best_image(data, 1920, 1440, data.get("fallback_keywords"))
    if img_bytes:
        sl.shapes.add_picture(_bytes_stream(img_bytes), 0, 0, W, H)
    else:
        keyword = data.get("keyword", data.get("title", ""))
        sl.shapes.add_picture(
            _styled_fallback_stream(1920, 1440, t["accent"], keyword), 0, 0, W, H)
    overlay_h = Inches(2.6)
    _rect(sl, 0, H - overlay_h, W, overlay_h, RGBColor(0x00, 0x00, 0x00))
    _rect(sl, 0, H - overlay_h - Inches(0.08), W, Inches(0.08), t["mid"])
    _textbox(sl, Inches(0.55), H - overlay_h + Inches(0.25),
             W - Inches(1.1), Inches(1.2), data.get("title", ""), 36, True, t["white"])
    body = data.get("body", "")
    if body:
        _textbox(sl, Inches(0.55), H - overlay_h + Inches(1.35),
                 W - Inches(1.1), Inches(0.8), body, 16, False, RGBColor(0xDD, 0xDD, 0xEE))
    caption = data.get("caption", "")
    if caption:
        _textbox(sl, Inches(0.55), H - Inches(0.35), W - Inches(1.1), Inches(0.25),
                 caption, 13, False, RGBColor(0xBB, 0xBB, 0xCC), italic=True)


def _build_image_text_slide(prs, data, t):
    """Photo left 45%, content right 55%. Uses Gemini pipeline."""
    W, H = prs.slide_width, prs.slide_height
    sl = prs.slides.add_slide(prs.slide_layouts[6])
    _set_bg(sl, t["bg"])
    hdr_h = Inches(1.25)
    _header_band(sl, W, t, hdr_h)
    _textbox(sl, Inches(0.55), Inches(0.18), W - Inches(1.1), hdr_h - Inches(0.22),
             data.get("title", ""), 26, True, t["white"])
    img_split = int(W * 0.44)
    img_top = hdr_h + Inches(0.12)
    img_h = H - img_top - Inches(0.22)
    img_bytes = _fetch_best_image(data, 960, 720, data.get("fallback_keywords"))
    if img_bytes:
        sl.shapes.add_picture(_bytes_stream(img_bytes),
                              Inches(0.12), img_top, img_split - Inches(0.2), img_h)
    else:
        keyword = data.get("keyword", data.get("title", ""))
        sl.shapes.add_picture(
            _styled_fallback_stream(960, 720, t["accent"], keyword),
            Inches(0.12), img_top, img_split - Inches(0.2), img_h)
    rx = img_split + Inches(0.22)
    rw = W - rx - Inches(0.35)
    ry = img_top + Inches(0.15)
    body = data.get("body", "")
    if body:
        _textbox(sl, rx, ry, rw, Inches(1.2), body, 14, False, t["muted"])
        ry += Inches(1.28)
    bullets = data.get("bullets", [])
    if bullets:
        tb = sl.shapes.add_textbox(rx, ry, rw, H - ry - Inches(0.4))
        tf = tb.text_frame
        tf.word_wrap = True
        for i, b in enumerate(bullets):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            p.text = f"\u25b8  {b}"
            p.space_before = Pt(8)
            for run in p.runs:
                run.font.size = Pt(16)
                run.font.color.rgb = t["dark"]
    _footer_rule(sl, W, H, t)


def _build_chart_slide(prs, data, t):
    W, H = prs.slide_width, prs.slide_height
    sl = prs.slides.add_slide(prs.slide_layouts[6])
    _set_bg(sl, t["bg"])
    hdr_h = Inches(1.25)
    _header_band(sl, W, t, hdr_h)
    _textbox(sl, Inches(0.55), Inches(0.18), W - Inches(1.1), hdr_h - Inches(0.22),
             data.get("title", "Chart"), 26, True, t["white"])
    chart_top = hdr_h + Inches(0.2)
    body = data.get("body", "")
    if body:
        _textbox(sl, Inches(0.6), chart_top, W - Inches(1.2), Inches(0.75),
                 body, 13, False, t["muted"])
        chart_top += Inches(0.82)
    labels = data.get("labels", [])
    values = data.get("values", [])
    chart_title = data.get("chart_title", data.get("title", ""))
    chart_type = data.get("chart_type", "bar").lower()
    accent_hex = t.get("mpl", "#1A4F8A")
    try:
        if chart_type == "line":
            buf = _chart_line(labels, values, chart_title, accent_hex)
        elif chart_type == "pie":
            buf = _chart_pie(labels, values, chart_title, accent_hex)
        elif chart_type == "horizontal_bar":
            buf = _chart_hbar(labels, values, chart_title, accent_hex)
        else:
            buf = _chart_bar(labels, values, chart_title, accent_hex)
        sl.shapes.add_picture(buf, Inches(0.5), chart_top,
                              W - Inches(1.0), H - chart_top - Inches(0.45))
    except Exception as exc:
        logger.error(f"Chart error: {exc}")
        _textbox(sl, Inches(0.6), chart_top + Inches(0.4), W - Inches(1.2),
                 Inches(1.0), f"[Chart error: {exc}]", 14, False, RGBColor(0x99, 0x00, 0x00))
    notes = data.get("notes", "")
    if notes:
        sl.notes_slide.notes_text_frame.text = notes
    _footer_rule(sl, W, H, t)


def _build_stat_cards_slide(prs, data, t):
    W, H = prs.slide_width, prs.slide_height
    sl = prs.slides.add_slide(prs.slide_layouts[6])
    _set_bg(sl, t["bg"])
    hdr_h = Inches(1.25)
    _header_band(sl, W, t, hdr_h)
    _textbox(sl, Inches(0.55), Inches(0.18), W - Inches(1.1), hdr_h - Inches(0.22),
             data.get("title", "Key Metrics"), 26, True, t["white"])
    card_top = hdr_h + Inches(0.32)
    body = data.get("body", "")
    if body:
        _textbox(sl, Inches(0.6), card_top, W - Inches(1.2), Inches(0.65),
                 body, 13, False, t["muted"])
        card_top += Inches(0.72)
    stats = data.get("stats", [])[:4]
    n = len(stats)
    if n == 0:
        _footer_rule(sl, W, H, t)
        return
    margin = Inches(0.5)
    gap = Inches(0.28)
    card_h = H - card_top - Inches(0.55)
    total_w = W - 2 * margin - gap * (n - 1)
    card_w = total_w // n
    trend_up   = RGBColor(0x0D, 0x7A, 0x3A)
    trend_down = RGBColor(0xBB, 0x1C, 0x1C)
    for i, stat in enumerate(stats):
        cx = margin + i * (card_w + gap)
        _rect(sl, cx, card_top, card_w, card_h, t["light"])
        _rect(sl, cx, card_top, card_w, Inches(0.14), t["accent"])
        _textbox(sl, cx + Inches(0.12), card_top + Inches(0.22),
                 card_w - Inches(0.24), Inches(1.55),
                 str(stat.get("value", "\u2013")), 50, True, t["accent"], align=PP_ALIGN.CENTER)
        trend = stat.get("trend", "")
        if trend == "up":
            _textbox(sl, cx + Inches(0.12), card_top + Inches(1.8),
                     card_w - Inches(0.24), Inches(0.4),
                     "\u25b2", 14, True, trend_up, align=PP_ALIGN.CENTER)
        elif trend == "down":
            _textbox(sl, cx + Inches(0.12), card_top + Inches(1.8),
                     card_w - Inches(0.24), Inches(0.4),
                     "\u25bc", 14, True, trend_down, align=PP_ALIGN.CENTER)
        label_y = card_top + Inches(1.82 if trend else 1.75)
        _textbox(sl, cx + Inches(0.1), label_y, card_w - Inches(0.2), Inches(0.65),
                 str(stat.get("label", "")), 15, True, t["dark"], align=PP_ALIGN.CENTER)
        detail = stat.get("detail", "")
        if detail:
            _textbox(sl, cx + Inches(0.1), label_y + Inches(0.68),
                     card_w - Inches(0.2), Inches(0.9),
                     detail, 11, False, t["muted"], align=PP_ALIGN.CENTER)
    _footer_rule(sl, W, H, t)


def _build_timeline_slide(prs, data, t):
    W, H = prs.slide_width, prs.slide_height
    sl = prs.slides.add_slide(prs.slide_layouts[6])
    _set_bg(sl, t["bg"])
    hdr_h = Inches(1.25)
    _header_band(sl, W, t, hdr_h)
    _textbox(sl, Inches(0.55), Inches(0.18), W - Inches(1.1), hdr_h - Inches(0.22),
             data.get("title", "Timeline"), 26, True, t["white"])
    body = data.get("body", "")
    if body:
        _textbox(sl, Inches(0.6), hdr_h + Inches(0.2), W - Inches(1.2), Inches(0.65),
                 body, 13, False, t["muted"])
    milestones = data.get("milestones", [])[:6]
    n = len(milestones)
    if n == 0:
        _footer_rule(sl, W, H, t)
        return
    margin = Inches(0.9)
    usable_w = W - 2 * margin
    step = usable_w // (n - 1) if n > 1 else usable_w // 2
    line_y = H // 2 + Inches(0.35)
    _rect(sl, margin, line_y - Inches(0.025), usable_w, Inches(0.05), t["accent"])
    dot_r = Inches(0.2)
    for i, ms in enumerate(milestones):
        x = margin + (i * step if n > 1 else usable_w // 2)
        above = (i % 2 == 0)
        _rect(sl, x - dot_r // 2, line_y - dot_r // 2, dot_r, dot_r, t["accent"])
        inner = dot_r // 2
        _rect(sl, x - inner // 2, line_y - inner // 2, inner, inner, t["white"])
        con_h = Inches(0.65)
        bw = Inches(1.65)
        bx = x - bw // 2
        if above:
            _rect(sl, x - Inches(0.015), line_y - dot_r // 2 - con_h,
                  Inches(0.03), con_h, t["mid"])
            _textbox(sl, bx, line_y - dot_r // 2 - con_h - Inches(0.5),
                     bw, Inches(0.4), str(ms.get("year", "")), 13, True, t["accent"],
                     align=PP_ALIGN.CENTER)
            _textbox(sl, bx, line_y - dot_r // 2 - con_h - Inches(0.95),
                     bw, Inches(0.42), str(ms.get("label", "")), 11, True, t["dark"],
                     align=PP_ALIGN.CENTER)
            if ms.get("detail"):
                _textbox(sl, bx - Inches(0.1), line_y - dot_r // 2 - con_h - Inches(1.42),
                         bw + Inches(0.2), Inches(0.44), ms["detail"], 9, False, t["muted"],
                         align=PP_ALIGN.CENTER)
        else:
            _rect(sl, x - Inches(0.015), line_y + dot_r // 2,
                  Inches(0.03), con_h, t["mid"])
            _textbox(sl, bx, line_y + dot_r // 2 + con_h + Inches(0.04),
                     bw, Inches(0.4), str(ms.get("year", "")), 13, True, t["accent"],
                     align=PP_ALIGN.CENTER)
            _textbox(sl, bx, line_y + dot_r // 2 + con_h + Inches(0.48),
                     bw, Inches(0.42), str(ms.get("label", "")), 11, True, t["dark"],
                     align=PP_ALIGN.CENTER)
            if ms.get("detail"):
                _textbox(sl, bx - Inches(0.1), line_y + dot_r // 2 + con_h + Inches(0.94),
                         bw + Inches(0.2), Inches(0.44), ms["detail"], 9, False, t["muted"],
                         align=PP_ALIGN.CENTER)
    _footer_rule(sl, W, H, t)


def _build_quote_slide(prs, data, t):
    W, H = prs.slide_width, prs.slide_height
    sl = prs.slides.add_slide(prs.slide_layouts[6])
    keyword = data.get("keyword", "")
    used_image = False
    if keyword:
        img_bytes = _fetch_best_image(data, 1920, 1080, data.get("fallback_keywords"))
        if img_bytes:
            sl.shapes.add_picture(_bytes_stream(img_bytes), 0, 0, W, H)
            _rect(sl, 0, 0, W, H, RGBColor(0x08, 0x08, 0x14))
            used_image = True
    if not used_image:
        _set_bg(sl, t["accent"])
    _textbox(sl, Inches(0.5), Inches(0.3), Inches(2.0), Inches(2.5),
             "\u201c", 120, True, t["mid"], align=PP_ALIGN.LEFT)
    _textbox(sl, Inches(1.1), Inches(1.55), W - Inches(2.2), Inches(3.4),
             data.get("quote", ""), 26, False, t["white"], align=PP_ALIGN.LEFT, italic=True)
    attribution = data.get("attribution", "")
    if attribution:
        _rect(sl, Inches(1.1), H - Inches(1.85), Inches(0.42), Inches(0.042), t["mid"])
        _textbox(sl, Inches(1.62), H - Inches(1.95), W - Inches(2.2), Inches(0.65),
                 f"\u2014 {attribution}", 15, True, t["light"])


def _build_section_slide(prs, data, t):
    W, H = prs.slide_width, prs.slide_height
    sl = prs.slides.add_slide(prs.slide_layouts[6])
    keyword = data.get("keyword", "")
    used_image = False
    if keyword:
        img_bytes = _fetch_best_image(data, 1920, 1080, data.get("fallback_keywords"))
        if img_bytes:
            sl.shapes.add_picture(_bytes_stream(img_bytes), 0, 0, W, H)
            _rect(sl, 0, 0, int(W * 0.52), H, t["accent"])
            used_image = True
    if not used_image:
        _set_bg(sl, t["bg"])
        _rect(sl, 0, 0, int(W * 0.52), H, t["accent"])
    num = str(data.get("section_number", ""))
    title_top = Inches(1.9) if not num else Inches(2.7)
    if num:
        _textbox(sl, Inches(0.55), Inches(0.9), Inches(4.5), Inches(1.6),
                 num, 88, True, t["mid"], align=PP_ALIGN.LEFT)
    _rect(sl, Inches(0.55), title_top - Inches(0.28),
          int(W * 0.36), Inches(0.055), t["mid"])
    _textbox(sl, Inches(0.55), title_top,
             int(W * 0.48) - Inches(0.7), Inches(2.6),
             data.get("title", "Section"), 34, True, t["white"])
    sub = data.get("subtitle", "")
    if sub:
        _textbox(sl, Inches(0.55), title_top + Inches(2.7),
                 int(W * 0.48) - Inches(0.7), Inches(0.85), sub, 17, False, t["light"])


def _build_agenda_slide(prs, data, t):
    W, H = prs.slide_width, prs.slide_height
    sl = prs.slides.add_slide(prs.slide_layouts[6])
    _set_bg(sl, t["bg"])
    hdr_h = Inches(1.25)
    _header_band(sl, W, t, hdr_h)
    _textbox(sl, Inches(0.55), Inches(0.18), W - Inches(1.1), hdr_h - Inches(0.22),
             data.get("title", "Agenda"), 26, True, t["white"])
    items = data.get("items", [])
    if not items:
        _footer_rule(sl, W, H, t)
        return
    two_col = len(items) > 4
    col_w = (W - Inches(1.2)) // 2 if two_col else W - Inches(1.2)
    avail_h = H - hdr_h - Inches(0.65)
    rows = (len(items) + 1) // 2 if two_col else len(items)
    row_h = min(Inches(0.95), avail_h / max(1, rows))
    for idx, item in enumerate(items):
        if two_col:
            col_idx = idx % 2
            row = idx // 2
            lx = Inches(0.6) if col_idx == 0 else Inches(0.6) + col_w + Inches(0.2)
        else:
            row = idx
            lx = Inches(0.6)
        ty = hdr_h + Inches(0.28) + row * row_h
        _rect(sl, lx, ty + Inches(0.06), Inches(0.44), Inches(0.44), t["accent"])
        _textbox(sl, lx, ty + Inches(0.05), Inches(0.44), Inches(0.44),
                 str(idx + 1), 15, True, t["white"], align=PP_ALIGN.CENTER)
        if isinstance(item, dict):
            label = item.get("label", "")
            desc  = item.get("description", "")
        else:
            label = str(item)
            desc  = ""
        _textbox(sl, lx + Inches(0.58), ty + Inches(0.08),
                 col_w - Inches(0.68), Inches(0.36), label, 16, True, t["dark"])
        if desc:
            _textbox(sl, lx + Inches(0.58), ty + Inches(0.46),
                     col_w - Inches(0.68), Inches(0.4), desc, 12, False, t["muted"])
    _footer_rule(sl, W, H, t)


# ---------------------------------------------------------------------------
# Builder registry
# ---------------------------------------------------------------------------

_BUILDERS = {
    "title":           _build_title_slide,
    "content":         _build_content_slide,
    "two_column":      _build_two_column_slide,
    "closing":         _build_closing_slide,
    "end":             _build_closing_slide,
    "thank_you":       _build_closing_slide,
    "image":           _build_image_slide,
    "image_text":      _build_image_text_slide,
    "chart":           _build_chart_slide,
    "stat_cards":      _build_stat_cards_slide,
    "timeline":        _build_timeline_slide,
    "quote":           _build_quote_slide,
    "section":         _build_section_slide,
    "section_divider": _build_section_slide,
    "agenda":          _build_agenda_slide,
}


# ---------------------------------------------------------------------------
# PptxToolset
# ---------------------------------------------------------------------------

class PptxToolset:
    """Enhanced toolset for generating PowerPoint presentations (v4)."""

    def __init__(self, host: str, port: int):
        self.host = host
        self.port = port
        self.output_dir = "outputs"
        os.makedirs(self.output_dir, exist_ok=True)
        logger.info(f"PptxToolset v4 ready. Output dir: ./{self.output_dir}")

    async def generate_pptx(
        self,
        filename: str,
        slides: list,
        theme: str = "blue",
    ) -> str:
        """
        Generates a PowerPoint (.pptx) from structured slide data.
        Images sourced via Gemini keyword generation + DuckDuckGo scraping,
        with Gemini Vision selecting the most contextually relevant candidate.
        """
        try:
            if isinstance(slides, str):
                slides = json.loads(slides)

            logger.info(f"Building '{filename}': {len(slides)} slides, theme={theme}")

            if not filename.endswith(".pptx"):
                filename += ".pptx"

            t = THEMES.get(theme.lower(), THEMES["blue"])
            prs = Presentation()
            prs.slide_width = Inches(13.33)
            prs.slide_height = Inches(7.5)

            for slide_data in slides:
                if not isinstance(slide_data, dict):
                    slide_data = {"type": "content", "title": str(slide_data)}
                slide_type = slide_data.get("type", "content").lower()
                builder = _BUILDERS.get(slide_type, _build_content_slide)
                builder(prs, slide_data, t)

            filepath = os.path.join(self.output_dir, filename)
            prs.save(filepath)
            url = f"http://{self.host}:{self.port}/outputs/{filename}"
            logger.info(f"Saved: {filepath}")
            return (
                f"\u2705 Successfully generated **{filename}**!\n\n"
                f"\U0001f4ca [Download your presentation here]({url})"
            )

        except Exception as exc:
            logger.error(f"PPTX generation error: {exc}", exc_info=True)
            return f"\u274c Failed to generate presentation: {exc}"

    def get_tools(self) -> dict[str, Any]:
        return {"generate_pptx": self}
