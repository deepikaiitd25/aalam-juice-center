"""
Enhanced PPTX Toolset v3
------------------------
Improvements over v2:
  - Richer content slides: body paragraph + bullets, better typography
  - New slide types: quote, section, agenda
  - 12 professional themes with complementary accent/bg/mid/muted palettes
  - Better title slide with optional full-bleed background image
  - Image fetching uses fallback_keywords for higher relevance
  - Multi-tone chart palettes with hue rotation
  - Cleaner visual hierarchy: gradient header bands, spacing, sizing
  - Stat cards support trend indicators (up/down arrows)
  - Timeline uses alternating above/below labels for readability

Stock photos: loremflickr.com (keyword-aware) → picsum.photos (fallback).
No API key required. No generative AI.
"""

from __future__ import annotations

import io
import json
import logging
import os
import urllib.request
from typing import Any

import matplotlib
matplotlib.use("Agg")
import matplotlib.pyplot as plt
import numpy as np
from PIL import Image

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
from pydantic import BaseModel

logger = logging.getLogger(__name__)


# ---------------------------------------------------------------------------
# Response model
# ---------------------------------------------------------------------------

class PptxGenerationResponse(BaseModel):
    status: str
    file_url: str | None = None
    error_message: str | None = None


# ---------------------------------------------------------------------------
# Themes – 12 professional palettes
# Keys: accent, mid, light, bg, white, dark, muted, mpl
# ---------------------------------------------------------------------------

THEMES: dict[str, dict] = {
    "blue": {
        "accent": RGBColor(0x1A, 0x4F, 0x8A),
        "mid":    RGBColor(0x26, 0x7A, 0xB5),
        "light":  RGBColor(0xE3, 0xEF, 0xF8),
        "bg":     RGBColor(0xFA, 0xFB, 0xFF),
        "white":  RGBColor(0xFF, 0xFF, 0xFF),
        "dark":   RGBColor(0x1A, 0x1A, 0x2E),
        "muted":  RGBColor(0x5C, 0x6B, 0x80),
        "mpl":    "#1A4F8A",
    },
    "green": {
        "accent": RGBColor(0x0D, 0x5C, 0x2E),
        "mid":    RGBColor(0x1D, 0x8A, 0x4A),
        "light":  RGBColor(0xD4, 0xED, 0xDA),
        "bg":     RGBColor(0xF6, 0xFB, 0xF7),
        "white":  RGBColor(0xFF, 0xFF, 0xFF),
        "dark":   RGBColor(0x0A, 0x1E, 0x12),
        "muted":  RGBColor(0x4A, 0x6A, 0x54),
        "mpl":    "#0D5C2E",
    },
    "dark": {
        "accent": RGBColor(0x0F, 0x1A, 0x35),
        "mid":    RGBColor(0x5B, 0x2D, 0x8E),
        "light":  RGBColor(0x2A, 0x2A, 0x45),
        "bg":     RGBColor(0x14, 0x14, 0x28),
        "white":  RGBColor(0xF0, 0xF0, 0xF8),
        "dark":   RGBColor(0xE8, 0xE8, 0xF5),
        "muted":  RGBColor(0x9A, 0x9A, 0xBB),
        "mpl":    "#5B2D8E",
    },
    "red": {
        "accent": RGBColor(0x8B, 0x14, 0x14),
        "mid":    RGBColor(0xC0, 0x2B, 0x2B),
        "light":  RGBColor(0xF9, 0xE0, 0xE0),
        "bg":     RGBColor(0xFF, 0xFA, 0xFA),
        "white":  RGBColor(0xFF, 0xFF, 0xFF),
        "dark":   RGBColor(0x1A, 0x08, 0x08),
        "muted":  RGBColor(0x7A, 0x55, 0x55),
        "mpl":    "#8B1414",
    },
    "purple": {
        "accent": RGBColor(0x4A, 0x0E, 0x7E),
        "mid":    RGBColor(0x7B, 0x2D, 0xBF),
        "light":  RGBColor(0xEA, 0xD9, 0xF7),
        "bg":     RGBColor(0xFB, 0xF7, 0xFF),
        "white":  RGBColor(0xFF, 0xFF, 0xFF),
        "dark":   RGBColor(0x18, 0x06, 0x2A),
        "muted":  RGBColor(0x6A, 0x55, 0x80),
        "mpl":    "#4A0E7E",
    },
    "orange": {
        "accent": RGBColor(0xC0, 0x4A, 0x08),
        "mid":    RGBColor(0xE8, 0x7A, 0x20),
        "light":  RGBColor(0xFD, 0xEA, 0xD5),
        "bg":     RGBColor(0xFF, 0xFB, 0xF7),
        "white":  RGBColor(0xFF, 0xFF, 0xFF),
        "dark":   RGBColor(0x1E, 0x0E, 0x00),
        "muted":  RGBColor(0x7A, 0x5A, 0x40),
        "mpl":    "#C04A08",
    },
    "teal": {
        "accent": RGBColor(0x02, 0x72, 0x80),
        "mid":    RGBColor(0x0A, 0xA8, 0xB8),
        "light":  RGBColor(0xCC, 0xF0, 0xF4),
        "bg":     RGBColor(0xF4, 0xFD, 0xFF),
        "white":  RGBColor(0xFF, 0xFF, 0xFF),
        "dark":   RGBColor(0x01, 0x20, 0x24),
        "muted":  RGBColor(0x40, 0x70, 0x76),
        "mpl":    "#027280",
    },
    "midnight": {
        "accent": RGBColor(0x0B, 0x24, 0x4F),
        "mid":    RGBColor(0xC9, 0xA0, 0x2A),
        "light":  RGBColor(0xE8, 0xEE, 0xF8),
        "bg":     RGBColor(0xF5, 0xF7, 0xFF),
        "white":  RGBColor(0xFF, 0xFF, 0xFF),
        "dark":   RGBColor(0x08, 0x14, 0x2E),
        "muted":  RGBColor(0x6A, 0x72, 0x90),
        "mpl":    "#0B244F",
    },
    "slate": {
        "accent": RGBColor(0x2E, 0x3A, 0x50),
        "mid":    RGBColor(0x48, 0xA9, 0xDC),
        "light":  RGBColor(0xE8, 0xEC, 0xF2),
        "bg":     RGBColor(0xF8, 0xF9, 0xFB),
        "white":  RGBColor(0xFF, 0xFF, 0xFF),
        "dark":   RGBColor(0x14, 0x1C, 0x2A),
        "muted":  RGBColor(0x7A, 0x85, 0x99),
        "mpl":    "#2E3A50",
    },
    "rose": {
        "accent": RGBColor(0x9C, 0x27, 0x50),
        "mid":    RGBColor(0xE0, 0x5C, 0x88),
        "light":  RGBColor(0xF9, 0xE0, 0xEB),
        "bg":     RGBColor(0xFF, 0xF7, 0xFA),
        "white":  RGBColor(0xFF, 0xFF, 0xFF),
        "dark":   RGBColor(0x28, 0x08, 0x16),
        "muted":  RGBColor(0x88, 0x5A, 0x6A),
        "mpl":    "#9C2750",
    },
    "forest": {
        "accent": RGBColor(0x1B, 0x3A, 0x24),
        "mid":    RGBColor(0xB8, 0x7A, 0x1E),
        "light":  RGBColor(0xD8, 0xE8, 0xD8),
        "bg":     RGBColor(0xF5, 0xFA, 0xF5),
        "white":  RGBColor(0xFF, 0xFF, 0xFF),
        "dark":   RGBColor(0x0A, 0x14, 0x0C),
        "muted":  RGBColor(0x5A, 0x6E, 0x5C),
        "mpl":    "#1B3A24",
    },
    "indigo": {
        "accent": RGBColor(0x31, 0x35, 0x9E),
        "mid":    RGBColor(0x60, 0x65, 0xDE),
        "light":  RGBColor(0xE2, 0xE3, 0xF8),
        "bg":     RGBColor(0xF7, 0xF7, 0xFF),
        "white":  RGBColor(0xFF, 0xFF, 0xFF),
        "dark":   RGBColor(0x0E, 0x10, 0x2E),
        "muted":  RGBColor(0x6A, 0x6C, 0xA8),
        "mpl":    "#31359E",
    },
}


# ---------------------------------------------------------------------------
# Image helpers
# ---------------------------------------------------------------------------

_IMG_TIMEOUT = 7


def _try_pexels(keyword: str, width: int, height: int) -> bytes | None:
    """Fetch high-quality image from Pexels (free, highly filtered for quality)."""
    # Enhance keyword specificity
    query_parts = keyword.split()
    # Remove generic words that might cause random matches
    filtered_parts = [p for p in query_parts if p.lower() not in ["a", "the", "and", "or"]]
    kw = urllib.request.quote(" ".join(filtered_parts) if filtered_parts else keyword, safe="")
    
    url = f"https://api.pexels.com/v1/search?query={kw}&per_page=5&size=large&sort=popular"
    try:
        req = urllib.request.Request(url, headers={
            "User-Agent": "pptx-agent/3.0",
            "Accept": "application/json"
        })
        with urllib.request.urlopen(req, timeout=_IMG_TIMEOUT) as resp:
            import json
            data = json.loads(resp.read().decode())
            photos = data.get("photos", [])
            # Filter for landscape orientation and reasonable size
            for photo in photos:
                if photo.get("width", 0) >= photo.get("height", 0) * 1.3:  # Landscape
                    photo_url = photo.get("src", {}).get("large")
                    if photo_url:
                        try:
                            img_req = urllib.request.Request(photo_url, headers={"User-Agent": "pptx-agent/3.0"})
                            with urllib.request.urlopen(img_req, timeout=_IMG_TIMEOUT) as img_resp:
                                img_data = img_resp.read()
                                if len(img_data) > 10_000:  # Quality threshold
                                    return img_data
                        except Exception:
                            pass
    except Exception as exc:
        pass
    return None


def _try_unsplash_raw(keyword: str, width: int, height: int) -> bytes | None:
    """Fetch from Unsplash source.unsplash.com - with specific topic filtering."""
    try:
        # Use more specific Unsplash topic parameter
        kw = urllib.request.quote(keyword, safe="")
        url = f"https://source.unsplash.com/{width}x{height}/?{kw},business,corporate"
        req = urllib.request.Request(url, headers={"User-Agent": "pptx-agent/3.0"})
        with urllib.request.urlopen(req, timeout=_IMG_TIMEOUT) as resp:
            img_data = resp.read()
            if len(img_data) > 10_000:
                return img_data
    except Exception as exc:
        pass
    return None


def _try_picsum(keyword: str, width: int, height: int) -> bytes | None:
    """Fallback to Picsum Photos with keyword-based seeding."""
    seed = abs(hash(keyword[:20])) % 5000
    url = f"https://picsum.photos/seed/{seed}/{width}/{height}"
    try:
        req = urllib.request.Request(url, headers={"User-Agent": "pptx-agent/3.0"})
        with urllib.request.urlopen(req, timeout=_IMG_TIMEOUT) as resp:
            data = resp.read()
        return data if len(data) > 5_000 else None
    except Exception as exc:
        pass
    return None


def _fetch_stock_photo(keyword: str, width: int = 1920, height: int = 1080,
                       fallback_keywords: list | None = None) -> bytes | None:
    """
    Fetch high-quality, contextually relevant images from multiple sources.
    Tries primary keyword first with each source before falling back.
    """
    keywords_to_try = [keyword]
    if fallback_keywords:
        keywords_to_try.extend(fallback_keywords)
    
    # Try each keyword with each source in order of quality
    for kw in keywords_to_try:
        # Pexels first (highest quality and most relevant)
        result = _try_pexels(kw, width, height)
        if result:
            return result
    
    # Second pass: Unsplash
    for kw in keywords_to_try:
        result = _try_unsplash_raw(kw, width, height)
        if result:
            return result
    
    # Final fallback: Picsum
    return _try_picsum(keyword, width, height)
    logger.warning(f"All image sources failed for '{keyword}'")
    return None


def _bytes_stream(data: bytes) -> io.BytesIO:
    buf = io.BytesIO(data)
    buf.seek(0)
    return buf


def _styled_fallback_stream(width_px: int, height_px: int,
                             accent: RGBColor, keyword: str = "") -> io.BytesIO:
    import colorsys
    r, g, b = accent.red / 255, accent.green / 255, accent.blue / 255
    h, s, v = colorsys.rgb_to_hsv(r, g, b)
    light_rgb = colorsys.hsv_to_rgb(h, max(0.10, s * 0.35), min(1.0, v * 1.25))

    fig_w, fig_h = width_px / 150, height_px / 150
    fig, ax = plt.subplots(figsize=(fig_w, fig_h), facecolor="white")
    ax.set_axis_off()

    grad = np.linspace(0, 1, 256).reshape(1, -1)
    grad = np.vstack([grad] * 64)
    ca, cb = [r, g, b], list(light_rgb)
    rgba = np.zeros((64, 256, 4))
    for ch in range(3):
        rgba[:, :, ch] = ca[ch] + (cb[ch] - ca[ch]) * grad
    rgba[:, :, 3] = 1.0
    ax.imshow(rgba, aspect="auto", extent=[0, 1, 0, 1], transform=ax.transAxes)
    if keyword:
        ax.text(0.5, 0.5, keyword.title(), transform=ax.transAxes,
                ha="center", va="center", fontsize=max(12, fig_w * 2),
                color="white", fontweight="bold", alpha=0.55, wrap=True)

    plt.tight_layout(pad=0)
    buf = io.BytesIO()
    fig.savefig(buf, format="PNG", dpi=150, bbox_inches="tight")
    plt.close(fig)
    buf.seek(0)
    return buf


# ---------------------------------------------------------------------------
# Chart builders
# ---------------------------------------------------------------------------

def _palette(accent_hex: str, n: int) -> list:
    import colorsys
    r = int(accent_hex[1:3], 16) / 255
    g = int(accent_hex[3:5], 16) / 255
    b = int(accent_hex[5:7], 16) / 255
    h, s, v = colorsys.rgb_to_hsv(r, g, b)
    out = []
    for i in range(n):
        ri, gi, bi = colorsys.hsv_to_rgb(
            (h + i * 0.08) % 1.0, max(0.25, s - i * 0.05), max(0.45, v - i * 0.04))
        out.append("#%02x%02x%02x" % (int(ri*255), int(gi*255), int(bi*255)))
    return out


def _chart_style(ax, title):
    ax.set_title(title, fontsize=14, fontweight="bold", color="#222", pad=14)
    ax.tick_params(axis="x", labelsize=10, colors="#555")
    ax.tick_params(axis="y", labelsize=9, colors="#888")
    ax.spines[["top", "right"]].set_visible(False)
    ax.spines[["left", "bottom"]].set_color("#DDD")
    ax.yaxis.grid(True, color="#F0F0F0", linewidth=0.8)
    ax.set_axisbelow(True)
    ax.set_facecolor("#FAFAFA")


def _chart_bar(labels, values, title, accent_hex, figsize=(10, 5)) -> io.BytesIO:
    colors = _palette(accent_hex, len(labels))
    fig, ax = plt.subplots(figsize=figsize, facecolor="white")
    bars = ax.bar(labels, values, color=colors, edgecolor="white", linewidth=1.0, width=0.6)
    _chart_style(ax, title)
    max_v = max(values) if values else 1
    for bar, val in zip(bars, values):
        label = f"{val:,}" if isinstance(val, int) else str(val)
        ax.text(bar.get_x() + bar.get_width() / 2, bar.get_height() + max_v * 0.012,
                label, ha="center", va="bottom", fontsize=9, color="#333", fontweight="bold")
    plt.tight_layout()
    buf = io.BytesIO()
    fig.savefig(buf, format="PNG", dpi=150)
    plt.close(fig)
    buf.seek(0)
    return buf


def _chart_line(labels, values, title, accent_hex, figsize=(10, 5)) -> io.BytesIO:
    fig, ax = plt.subplots(figsize=figsize, facecolor="white")
    ax.plot(labels, values, color=accent_hex, linewidth=2.5, marker="o",
            markersize=8, markerfacecolor="white",
            markeredgewidth=2.5, markeredgecolor=accent_hex)
    ax.fill_between(range(len(labels)), values, alpha=0.10, color=accent_hex)
    ax.set_xticks(range(len(labels)))
    ax.set_xticklabels(labels, fontsize=10, color="#555")
    _chart_style(ax, title)
    plt.tight_layout()
    buf = io.BytesIO()
    fig.savefig(buf, format="PNG", dpi=150)
    plt.close(fig)
    buf.seek(0)
    return buf


def _chart_pie(labels, values, title, accent_hex, figsize=(7, 4.5)) -> io.BytesIO:
    colors = _palette(accent_hex, len(labels))
    fig, ax = plt.subplots(figsize=figsize, facecolor="white")
    wedges, _, autotexts = ax.pie(
        values, labels=None, colors=colors,
        autopct="%1.1f%%", startangle=140,
        wedgeprops=dict(width=0.6, edgecolor="white", linewidth=2),
        pctdistance=0.77,
    )
    for at in autotexts:
        at.set_fontsize(9)
        at.set_color("white")
        at.set_fontweight("bold")
    ax.legend(wedges, labels, loc="center left",
              bbox_to_anchor=(1, 0, 0.5, 1), fontsize=9, frameon=False)
    ax.set_title(title, fontsize=14, fontweight="bold", color="#222", pad=14)
    plt.tight_layout()
    buf = io.BytesIO()
    fig.savefig(buf, format="PNG", dpi=150, bbox_inches="tight")
    plt.close(fig)
    buf.seek(0)
    return buf


def _chart_hbar(labels, values, title, accent_hex, figsize=(9, 5)) -> io.BytesIO:
    colors = _palette(accent_hex, len(labels))
    fig, ax = plt.subplots(figsize=figsize, facecolor="white")
    bars = ax.barh(labels[::-1], values[::-1], color=colors[::-1],
                   edgecolor="white", linewidth=0.8, height=0.55)
    ax.set_title(title, fontsize=14, fontweight="bold", color="#222", pad=12)
    ax.tick_params(axis="y", labelsize=11, colors="#444")
    ax.tick_params(axis="x", labelsize=9, colors="#888")
    ax.spines[["top", "right", "left"]].set_visible(False)
    ax.spines["bottom"].set_color("#DDD")
    ax.xaxis.grid(True, color="#F0F0F0", linewidth=0.8)
    ax.set_axisbelow(True)
    ax.set_facecolor("#FAFAFA")
    max_v = max(values) if values else 1
    for bar, val in zip(bars, values[::-1]):
        label = f"{val:,}" if isinstance(val, int) else str(val)
        ax.text(val + max_v * 0.01, bar.get_y() + bar.get_height() / 2,
                label, va="center", fontsize=9, color="#333", fontweight="bold")
    plt.tight_layout()
    buf = io.BytesIO()
    fig.savefig(buf, format="PNG", dpi=150)
    plt.close(fig)
    buf.seek(0)
    return buf


# ---------------------------------------------------------------------------
# Low-level drawing helpers
# ---------------------------------------------------------------------------

def _set_bg(slide, color: RGBColor):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def _rect(slide, left, top, width, height, color: RGBColor):
    shp = slide.shapes.add_shape(1, left, top, width, height)
    shp.fill.solid()
    shp.fill.fore_color.rgb = color
    shp.line.fill.background()
    return shp


def _textbox(slide, left, top, width, height, text, size, bold, color: RGBColor,
             align=PP_ALIGN.LEFT, wrap=True, italic=False):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = wrap
    tf.text = text
    for para in tf.paragraphs:
        para.alignment = align
        for run in para.runs:
            run.font.size = Pt(size)
            run.font.bold = bold
            run.font.italic = italic
            run.font.color.rgb = color
    return tb


def _header_band(slide, W, t, height=None):
    """Full-width accent header band with a mid-tone bottom stripe."""
    if height is None:
        height = Inches(1.25)
    _rect(slide, 0, 0, W, height, t["accent"])
    _rect(slide, 0, height - Inches(0.055), W, Inches(0.055), t["mid"])


def _footer_rule(slide, W, H, t):
    _rect(slide, 0, H - Inches(0.075), W, Inches(0.075), t["accent"])


# ---------------------------------------------------------------------------
# Slide builders
# ---------------------------------------------------------------------------

def _build_title_slide(prs, data, t):
    """
    Cover slide. Optional background image via 'keyword'.
    Fields: title, subtitle, keyword (optional), fallback_keywords (optional list),
            presenter (optional)
    """
    W, H = prs.slide_width, prs.slide_height
    sl = prs.slides.add_slide(prs.slide_layouts[6])

    keyword = data.get("keyword", "")
    used_image = False
    if keyword:
        img_bytes = _fetch_stock_photo(keyword, 1920, 1080,
                                       data.get("fallback_keywords"))
        if img_bytes:
            sl.shapes.add_picture(_bytes_stream(img_bytes), 0, 0, W, H)
            # Dark overlay covering left 55%
            _rect(sl, 0, 0, int(W * 0.55), H, RGBColor(0x08, 0x08, 0x10))
            # Softer overlay on right
            _rect(sl, int(W * 0.45), 0, int(W * 0.55), H, RGBColor(0x08, 0x08, 0x14))
            # Accent stripe
            _rect(sl, 0, 0, Inches(0.35), H, t["mid"])
            used_image = True

    if not used_image:
        _set_bg(sl, t["accent"])
        _rect(sl, 0, 0, Inches(0.4), H, t["mid"])
        _rect(sl, 0, H - Inches(1.9), W, Inches(1.9), t["bg"])

    text_col = t["white"]
    sub_col = t["light"]

    _textbox(sl, Inches(0.85), Inches(1.1), W - Inches(1.7), Inches(2.8),
             data.get("title", "Presentation"), 42, True, text_col)
    sub = data.get("subtitle", "")
    if sub:
        _textbox(sl, Inches(0.85), Inches(4.1), W - Inches(1.7), Inches(1.1),
                 sub, 21, False, sub_col)
    presenter = data.get("presenter", "")
    if presenter:
        _textbox(sl, Inches(0.85), H - Inches(1.05), W - Inches(1.7), Inches(0.65),
                 presenter, 14, False,
                 RGBColor(0xBB, 0xBB, 0xCC) if used_image else t["muted"],
                 italic=True)


def _build_content_slide(prs, data, t):
    """
    Content slide with optional intro paragraph + bullet list.
    Fields: title, body (optional paragraph), bullets (list[str]), notes
    """
    W, H = prs.slide_width, prs.slide_height
    sl = prs.slides.add_slide(prs.slide_layouts[6])
    _set_bg(sl, t["bg"])

    hdr_h = Inches(1.25)
    _header_band(sl, W, t, hdr_h)
    _textbox(sl, Inches(0.55), Inches(0.18), W - Inches(1.1), hdr_h - Inches(0.22),
             data.get("title", "Slide"), 26, True, t["white"])

    y = hdr_h + Inches(0.28)
    body = data.get("body", "")
    if body:
        _textbox(sl, Inches(0.6), y, W - Inches(1.2), Inches(1.25),
                 body, 15, False, t["dark"])
        y += Inches(1.3)

    bullets = data.get("bullets", [])
    if bullets:
        tb = sl.shapes.add_textbox(Inches(0.6), y, W - Inches(1.2),
                                   H - y - Inches(0.55))
        tf = tb.text_frame
        tf.word_wrap = True
        for i, b in enumerate(bullets):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            p.text = f"\u25b8  {b}"
            p.space_before = Pt(8)
            for run in p.runs:
                run.font.size = Pt(17)
                run.font.color.rgb = t["dark"]

    notes = data.get("notes", "")
    if notes:
        sl.notes_slide.notes_text_frame.text = notes
    _footer_rule(sl, W, H, t)


def _build_two_column_slide(prs, data, t):
    """
    Side-by-side comparison. Each column supports body + bullets.
    Fields: title, left_title, left_body, left_bullets,
                    right_title, right_body, right_bullets
    """
    W, H = prs.slide_width, prs.slide_height
    sl = prs.slides.add_slide(prs.slide_layouts[6])
    _set_bg(sl, t["bg"])

    hdr_h = Inches(1.25)
    _header_band(sl, W, t, hdr_h)
    _textbox(sl, Inches(0.55), Inches(0.18), W - Inches(1.1), hdr_h - Inches(0.22),
             data.get("title", "Comparison"), 26, True, t["white"])

    mid = W // 2
    col_top = hdr_h + Inches(0.22)
    col_h = H - col_top - Inches(0.55)
    col_w = mid - Inches(0.85)
    _rect(sl, mid - Inches(0.025), col_top, Inches(0.05), col_h, t["light"])

    def _col(lx, ctitle, cbody, cbullets):
        _rect(sl, lx, col_top, Inches(0.07), Inches(0.5), t["mid"])
        _textbox(sl, lx + Inches(0.15), col_top + Inches(0.04),
                 col_w - Inches(0.12), Inches(0.48),
                 ctitle, 18, True, t["accent"])
        cy = col_top + Inches(0.6)
        if cbody:
            _textbox(sl, lx + Inches(0.1), cy, col_w - Inches(0.1), Inches(1.0),
                     cbody, 13, False, t["muted"])
            cy += Inches(1.05)
        if cbullets:
            tb = sl.shapes.add_textbox(lx + Inches(0.1), cy,
                                       col_w - Inches(0.1), col_h - (cy - col_top))
            tf = tb.text_frame
            tf.word_wrap = True
            for i, b in enumerate(cbullets):
                p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
                p.text = f"\u25b8  {b}"
                p.space_before = Pt(7)
                for run in p.runs:
                    run.font.size = Pt(15)
                    run.font.color.rgb = t["dark"]

    _col(Inches(0.5), data.get("left_title", "Option A"),
         data.get("left_body", ""), data.get("left_bullets", []))
    _col(mid + Inches(0.35), data.get("right_title", "Option B"),
         data.get("right_body", ""), data.get("right_bullets", []))
    _footer_rule(sl, W, H, t)


def _build_closing_slide(prs, data, t):
    """Closing/thank-you slide. Fields: title, subtitle, contact"""
    W, H = prs.slide_width, prs.slide_height
    sl = prs.slides.add_slide(prs.slide_layouts[6])
    _set_bg(sl, t["accent"])
    _rect(sl, 0, 0, Inches(0.45), H, t["mid"])
    _rect(sl, W - Inches(3.8), H - Inches(3.8), Inches(3.8), Inches(3.8), t["mid"])
    _rect(sl, 0, H - Inches(2.1), W, Inches(2.1), t["mid"])
    _textbox(sl, Inches(1.3), Inches(1.5), W - Inches(2.2), Inches(2.9),
             data.get("title", "Thank You"), 50, True, t["white"],
             align=PP_ALIGN.CENTER)
    sub = data.get("subtitle", "")
    if sub:
        _textbox(sl, Inches(1.3), Inches(4.6), W - Inches(2.2), Inches(0.95),
                 sub, 20, False, t["light"], align=PP_ALIGN.CENTER)
    contact = data.get("contact", "")
    if contact:
        _textbox(sl, Inches(1.3), H - Inches(1.65), W - Inches(2.2), Inches(0.65),
                 contact, 13, False, t["light"], align=PP_ALIGN.CENTER, italic=True)


def _build_image_slide(prs, data, t):
    """
    Full-bleed photo + title overlay with improved formatting.
    Fields: title, keyword, fallback_keywords, caption, body
    """
    W, H = prs.slide_width, prs.slide_height
    sl = prs.slides.add_slide(prs.slide_layouts[6])

    keyword = data.get("keyword", data.get("title", "corporate"))
    fallbacks = data.get("fallback_keywords", [])
    img_bytes = _fetch_stock_photo(keyword, 1920, 1440, fallbacks)

    if img_bytes:
        sl.shapes.add_picture(_bytes_stream(img_bytes), 0, 0, W, H)
    else:
        sl.shapes.add_picture(_styled_fallback_stream(1920, 1440, t["accent"], keyword),
                              0, 0, W, H)

    # Enhanced dark overlay for better text readability
    overlay_h = Inches(2.6)
    _rect(sl, 0, H - overlay_h, W, overlay_h, RGBColor(0x00, 0x00, 0x00))
    _rect(sl, 0, H - overlay_h - Inches(0.08), W, Inches(0.08), t["mid"])

    # Title with better spacing
    _textbox(sl, Inches(0.55), H - overlay_h + Inches(0.25),
             W - Inches(1.1), Inches(1.2),
             data.get("title", ""), 36, True, t["white"])
    
    # Body text if provided
    body = data.get("body", "")
    if body:
        _textbox(sl, Inches(0.55), H - overlay_h + Inches(1.35),
                 W - Inches(1.1), Inches(0.8),
                 body, 16, False, RGBColor(0xDD, 0xDD, 0xEE))
    
    caption = data.get("caption", "")
    if caption:
        _textbox(sl, Inches(0.55), H - Inches(0.35), W - Inches(1.1), Inches(0.25),
                 caption, 13, False, RGBColor(0xBB, 0xBB, 0xCC), italic=True)


def _build_image_text_slide(prs, data, t):
    """
    Photo left 45%, content right 55%.
    Fields: title, keyword, fallback_keywords, body, bullets
    """
    W, H = prs.slide_width, prs.slide_height
    sl = prs.slides.add_slide(prs.slide_layouts[6])
    _set_bg(sl, t["bg"])

    hdr_h = Inches(1.25)
    _header_band(sl, W, t, hdr_h)
    _textbox(sl, Inches(0.55), Inches(0.18), W - Inches(1.1), hdr_h - Inches(0.22),
             data.get("title", ""), 26, True, t["white"])

    img_split = int(W * 0.44)
    img_top = hdr_h + Inches(0.12)
    img_h = H - img_top - Inches(0.22)

    keyword = data.get("keyword", data.get("title", "business"))
    img_bytes = _fetch_stock_photo(keyword, 960, 720, data.get("fallback_keywords"))
    if img_bytes:
        sl.shapes.add_picture(_bytes_stream(img_bytes),
                              Inches(0.12), img_top, img_split - Inches(0.2), img_h)
    else:
        sl.shapes.add_picture(_styled_fallback_stream(960, 720, t["accent"], keyword),
                              Inches(0.12), img_top, img_split - Inches(0.2), img_h)

    rx = img_split + Inches(0.22)
    rw = W - rx - Inches(0.35)
    ry = img_top + Inches(0.15)

    body = data.get("body", "")
    if body:
        _textbox(sl, rx, ry, rw, Inches(1.2), body, 14, False, t["muted"])
        ry += Inches(1.28)

    bullets = data.get("bullets", [])
    if bullets:
        tb = sl.shapes.add_textbox(rx, ry, rw, H - ry - Inches(0.4))
        tf = tb.text_frame
        tf.word_wrap = True
        for i, b in enumerate(bullets):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            p.text = f"\u25b8  {b}"
            p.space_before = Pt(8)
            for run in p.runs:
                run.font.size = Pt(16)
                run.font.color.rgb = t["dark"]

    _footer_rule(sl, W, H, t)


def _build_chart_slide(prs, data, t):
    """
    Matplotlib chart.
    Fields: title, body, chart_type, chart_title, labels, values, notes
    """
    W, H = prs.slide_width, prs.slide_height
    sl = prs.slides.add_slide(prs.slide_layouts[6])
    _set_bg(sl, t["bg"])

    hdr_h = Inches(1.25)
    _header_band(sl, W, t, hdr_h)
    _textbox(sl, Inches(0.55), Inches(0.18), W - Inches(1.1), hdr_h - Inches(0.22),
             data.get("title", "Chart"), 26, True, t["white"])

    chart_top = hdr_h + Inches(0.2)
    body = data.get("body", "")
    if body:
        _textbox(sl, Inches(0.6), chart_top, W - Inches(1.2), Inches(0.75),
                 body, 13, False, t["muted"])
        chart_top += Inches(0.82)

    labels = data.get("labels", [])
    values = data.get("values", [])
    chart_title = data.get("chart_title", data.get("title", ""))
    chart_type = data.get("chart_type", "bar").lower()
    accent_hex = t.get("mpl", "#1A4F8A")

    try:
        if chart_type == "line":
            buf = _chart_line(labels, values, chart_title, accent_hex)
        elif chart_type == "pie":
            buf = _chart_pie(labels, values, chart_title, accent_hex)
        elif chart_type == "horizontal_bar":
            buf = _chart_hbar(labels, values, chart_title, accent_hex)
        else:
            buf = _chart_bar(labels, values, chart_title, accent_hex)
        sl.shapes.add_picture(buf, Inches(0.5), chart_top,
                              W - Inches(1.0), H - chart_top - Inches(0.45))
    except Exception as exc:
        logger.error(f"Chart error: {exc}")
        _textbox(sl, Inches(0.6), chart_top + Inches(0.4), W - Inches(1.2),
                 Inches(1.0), f"[Chart error: {exc}]", 14, False,
                 RGBColor(0x99, 0x00, 0x00))

    notes = data.get("notes", "")
    if notes:
        sl.notes_slide.notes_text_frame.text = notes
    _footer_rule(sl, W, H, t)


def _build_stat_cards_slide(prs, data, t):
    """
    2–4 KPI stat cards with optional trend indicators.
    Fields: title, body, stats [{value, label, detail, trend: "up"|"down"|""}]
    """
    W, H = prs.slide_width, prs.slide_height
    sl = prs.slides.add_slide(prs.slide_layouts[6])
    _set_bg(sl, t["bg"])

    hdr_h = Inches(1.25)
    _header_band(sl, W, t, hdr_h)
    _textbox(sl, Inches(0.55), Inches(0.18), W - Inches(1.1), hdr_h - Inches(0.22),
             data.get("title", "Key Metrics"), 26, True, t["white"])

    card_top = hdr_h + Inches(0.32)
    body = data.get("body", "")
    if body:
        _textbox(sl, Inches(0.6), card_top, W - Inches(1.2), Inches(0.65),
                 body, 13, False, t["muted"])
        card_top += Inches(0.72)

    stats = data.get("stats", [])[:4]
    n = len(stats)
    if n == 0:
        _footer_rule(sl, W, H, t)
        return

    margin = Inches(0.5)
    gap = Inches(0.28)
    card_h = H - card_top - Inches(0.55)
    total_w = W - 2 * margin - gap * (n - 1)
    card_w = total_w // n

    trend_up   = RGBColor(0x0D, 0x7A, 0x3A)
    trend_down = RGBColor(0xBB, 0x1C, 0x1C)

    for i, stat in enumerate(stats):
        cx = margin + i * (card_w + gap)
        _rect(sl, cx, card_top, card_w, card_h, t["light"])
        _rect(sl, cx, card_top, card_w, Inches(0.14), t["accent"])

        _textbox(sl, cx + Inches(0.12), card_top + Inches(0.22),
                 card_w - Inches(0.24), Inches(1.55),
                 str(stat.get("value", "\u2013")),
                 50, True, t["accent"], align=PP_ALIGN.CENTER)

        trend = stat.get("trend", "")
        if trend == "up":
            _textbox(sl, cx + Inches(0.12), card_top + Inches(1.8),
                     card_w - Inches(0.24), Inches(0.4),
                     "\u25b2", 14, True, trend_up, align=PP_ALIGN.CENTER)
        elif trend == "down":
            _textbox(sl, cx + Inches(0.12), card_top + Inches(1.8),
                     card_w - Inches(0.24), Inches(0.4),
                     "\u25bc", 14, True, trend_down, align=PP_ALIGN.CENTER)

        label_y = card_top + Inches(1.82 if trend else 1.75)
        _textbox(sl, cx + Inches(0.1), label_y,
                 card_w - Inches(0.2), Inches(0.65),
                 str(stat.get("label", "")),
                 15, True, t["dark"], align=PP_ALIGN.CENTER)

        detail = stat.get("detail", "")
        if detail:
            _textbox(sl, cx + Inches(0.1), label_y + Inches(0.68),
                     card_w - Inches(0.2), Inches(0.9),
                     detail, 11, False, t["muted"], align=PP_ALIGN.CENTER)

    _footer_rule(sl, W, H, t)


def _build_timeline_slide(prs, data, t):
    """
    Horizontal milestone timeline, alternating above/below labels.
    Fields: title, body, milestones [{year, label, detail}]
    """
    W, H = prs.slide_width, prs.slide_height
    sl = prs.slides.add_slide(prs.slide_layouts[6])
    _set_bg(sl, t["bg"])

    hdr_h = Inches(1.25)
    _header_band(sl, W, t, hdr_h)
    _textbox(sl, Inches(0.55), Inches(0.18), W - Inches(1.1), hdr_h - Inches(0.22),
             data.get("title", "Timeline"), 26, True, t["white"])

    body = data.get("body", "")
    if body:
        _textbox(sl, Inches(0.6), hdr_h + Inches(0.2), W - Inches(1.2), Inches(0.65),
                 body, 13, False, t["muted"])

    milestones = data.get("milestones", [])[:6]
    n = len(milestones)
    if n == 0:
        _footer_rule(sl, W, H, t)
        return

    margin = Inches(0.9)
    usable_w = W - 2 * margin
    step = usable_w // (n - 1) if n > 1 else usable_w // 2
    line_y = H // 2 + Inches(0.35)

    # Spine
    _rect(sl, margin, line_y - Inches(0.025), usable_w, Inches(0.05), t["accent"])

    dot_r = Inches(0.2)
    for i, ms in enumerate(milestones):
        x = margin + (i * step if n > 1 else usable_w // 2)
        above = (i % 2 == 0)

        _rect(sl, x - dot_r // 2, line_y - dot_r // 2, dot_r, dot_r, t["accent"])
        inner = dot_r // 2
        _rect(sl, x - inner // 2, line_y - inner // 2, inner, inner, t["white"])

        con_h = Inches(0.65)
        bw = Inches(1.65)
        bx = x - bw // 2

        if above:
            _rect(sl, x - Inches(0.015), line_y - dot_r // 2 - con_h,
                  Inches(0.03), con_h, t["mid"])
            _textbox(sl, bx, line_y - dot_r // 2 - con_h - Inches(0.5),
                     bw, Inches(0.4),
                     str(ms.get("year", "")), 13, True, t["accent"],
                     align=PP_ALIGN.CENTER)
            _textbox(sl, bx, line_y - dot_r // 2 - con_h - Inches(0.95),
                     bw, Inches(0.42),
                     str(ms.get("label", "")), 11, True, t["dark"],
                     align=PP_ALIGN.CENTER)
            if ms.get("detail"):
                _textbox(sl, bx - Inches(0.1), line_y - dot_r // 2 - con_h - Inches(1.42),
                         bw + Inches(0.2), Inches(0.44),
                         ms["detail"], 9, False, t["muted"], align=PP_ALIGN.CENTER)
        else:
            _rect(sl, x - Inches(0.015), line_y + dot_r // 2,
                  Inches(0.03), con_h, t["mid"])
            _textbox(sl, bx, line_y + dot_r // 2 + con_h + Inches(0.04),
                     bw, Inches(0.4),
                     str(ms.get("year", "")), 13, True, t["accent"],
                     align=PP_ALIGN.CENTER)
            _textbox(sl, bx, line_y + dot_r // 2 + con_h + Inches(0.48),
                     bw, Inches(0.42),
                     str(ms.get("label", "")), 11, True, t["dark"],
                     align=PP_ALIGN.CENTER)
            if ms.get("detail"):
                _textbox(sl, bx - Inches(0.1), line_y + dot_r // 2 + con_h + Inches(0.94),
                         bw + Inches(0.2), Inches(0.44),
                         ms["detail"], 9, False, t["muted"], align=PP_ALIGN.CENTER)

    _footer_rule(sl, W, H, t)


# ── New in v3 ──────────────────────────────────────────────────────────────

def _build_quote_slide(prs, data, t):
    """
    Large pull-quote slide.
    Fields: quote, attribution, keyword (optional background image), fallback_keywords
    """
    W, H = prs.slide_width, prs.slide_height
    sl = prs.slides.add_slide(prs.slide_layouts[6])

    keyword = data.get("keyword", "")
    used_image = False
    if keyword:
        img_bytes = _fetch_stock_photo(keyword, 1920, 1080, data.get("fallback_keywords"))
        if img_bytes:
            sl.shapes.add_picture(_bytes_stream(img_bytes), 0, 0, W, H)
            _rect(sl, 0, 0, W, H, RGBColor(0x08, 0x08, 0x14))
            used_image = True

    if not used_image:
        _set_bg(sl, t["accent"])
        _rect(sl, 0, 0, W, H,
              RGBColor(max(0, t["accent"].red - 8),
                       max(0, t["accent"].green - 8),
                       max(0, t["accent"].blue - 8)))

    # Big decorative opening quote mark
    _textbox(sl, Inches(0.5), Inches(0.3), Inches(2.0), Inches(2.5),
             "\u201c", 120, True, t["mid"], align=PP_ALIGN.LEFT)

    _textbox(sl, Inches(1.1), Inches(1.55), W - Inches(2.2), Inches(3.4),
             data.get("quote", ""), 26, False, t["white"],
             align=PP_ALIGN.LEFT, italic=True)

    attribution = data.get("attribution", "")
    if attribution:
        _rect(sl, Inches(1.1), H - Inches(1.85), Inches(0.42), Inches(0.042), t["mid"])
        _textbox(sl, Inches(1.62), H - Inches(1.95), W - Inches(2.2), Inches(0.65),
                 f"\u2014 {attribution}", 15, True, t["light"])


def _build_section_slide(prs, data, t):
    """
    Bold section divider slide.
    Fields: section_number (optional), title, subtitle, keyword (optional), fallback_keywords
    """
    W, H = prs.slide_width, prs.slide_height
    sl = prs.slides.add_slide(prs.slide_layouts[6])

    keyword = data.get("keyword", "")
    used_image = False
    if keyword:
        img_bytes = _fetch_stock_photo(keyword, 1920, 1080, data.get("fallback_keywords"))
        if img_bytes:
            sl.shapes.add_picture(_bytes_stream(img_bytes), 0, 0, W, H)
            _rect(sl, 0, 0, int(W * 0.52), H, t["accent"])
            used_image = True

    if not used_image:
        _set_bg(sl, t["bg"])
        _rect(sl, 0, 0, int(W * 0.52), H, t["accent"])

    num = str(data.get("section_number", ""))
    title_top = Inches(1.9) if not num else Inches(2.7)

    if num:
        _textbox(sl, Inches(0.55), Inches(0.9), Inches(4.5), Inches(1.6),
                 num, 88, True, t["mid"], align=PP_ALIGN.LEFT)

    _rect(sl, Inches(0.55), title_top - Inches(0.28),
          int(W * 0.36), Inches(0.055), t["mid"])

    _textbox(sl, Inches(0.55), title_top,
             int(W * 0.48) - Inches(0.7), Inches(2.6),
             data.get("title", "Section"), 34, True, t["white"])

    sub = data.get("subtitle", "")
    if sub:
        _textbox(sl, Inches(0.55), title_top + Inches(2.7),
                 int(W * 0.48) - Inches(0.7), Inches(0.85),
                 sub, 17, False, t["light"])


def _build_agenda_slide(prs, data, t):
    """
    Numbered agenda / outline slide.
    Fields: title, items (list[str] or list[{label, description}])
    """
    W, H = prs.slide_width, prs.slide_height
    sl = prs.slides.add_slide(prs.slide_layouts[6])
    _set_bg(sl, t["bg"])

    hdr_h = Inches(1.25)
    _header_band(sl, W, t, hdr_h)
    _textbox(sl, Inches(0.55), Inches(0.18), W - Inches(1.1), hdr_h - Inches(0.22),
             data.get("title", "Agenda"), 26, True, t["white"])

    items = data.get("items", [])
    if not items:
        _footer_rule(sl, W, H, t)
        return

    two_col = len(items) > 4
    col_w = (W - Inches(1.2)) // 2 if two_col else W - Inches(1.2)
    avail_h = H - hdr_h - Inches(0.65)
    rows = (len(items) + 1) // 2 if two_col else len(items)
    row_h = min(Inches(0.95), avail_h / max(1, rows))

    for idx, item in enumerate(items):
        if two_col:
            col_idx = idx % 2
            row = idx // 2
            lx = Inches(0.6) if col_idx == 0 else Inches(0.6) + col_w + Inches(0.2)
        else:
            row = idx
            lx = Inches(0.6)

        ty = hdr_h + Inches(0.28) + row * row_h
        _rect(sl, lx, ty + Inches(0.06), Inches(0.44), Inches(0.44), t["accent"])
        _textbox(sl, lx, ty + Inches(0.05), Inches(0.44), Inches(0.44),
                 str(idx + 1), 15, True, t["white"], align=PP_ALIGN.CENTER)

        if isinstance(item, dict):
            label = item.get("label", "")
            desc = item.get("description", "")
        else:
            label = str(item)
            desc = ""

        _textbox(sl, lx + Inches(0.58), ty + Inches(0.08),
                 col_w - Inches(0.68), Inches(0.36),
                 label, 16, True, t["dark"])
        if desc:
            _textbox(sl, lx + Inches(0.58), ty + Inches(0.46),
                     col_w - Inches(0.68), Inches(0.4),
                     desc, 12, False, t["muted"])

    _footer_rule(sl, W, H, t)


# ---------------------------------------------------------------------------
# Builder registry
# ---------------------------------------------------------------------------

_BUILDERS = {
    "title":           _build_title_slide,
    "content":         _build_content_slide,
    "two_column":      _build_two_column_slide,
    "closing":         _build_closing_slide,
    "end":             _build_closing_slide,
    "thank_you":       _build_closing_slide,
    "image":           _build_image_slide,
    "image_text":      _build_image_text_slide,
    "chart":           _build_chart_slide,
    "stat_cards":      _build_stat_cards_slide,
    "timeline":        _build_timeline_slide,
    # v3 new
    "quote":           _build_quote_slide,
    "section":         _build_section_slide,
    "section_divider": _build_section_slide,
    "agenda":          _build_agenda_slide,
}


# ---------------------------------------------------------------------------
# PptxToolset
# ---------------------------------------------------------------------------

class PptxToolset:
    """Enhanced toolset for generating PowerPoint presentations (v3)."""

    def __init__(self, host: str, port: int):
        self.host = host
        self.port = port
        self.output_dir = "outputs"
        os.makedirs(self.output_dir, exist_ok=True)
        logger.info(f"PptxToolset v3 ready. Output dir: ./{self.output_dir}")

    async def generate_pptx(
        self,
        filename: str,
        slides: list,
        theme: str = "blue",
    ) -> str:
        """
        Generates a PowerPoint (.pptx) from structured slide data.

        Args:
            filename:  Output filename without extension, e.g. "startup_pitch"
            slides:    List of slide descriptor dicts; each must have a "type" field.
            theme:     Color theme name — see THEMES dict for all 12 options.
        """
        try:
            if isinstance(slides, str):
                slides = json.loads(slides)

            logger.info(f"Building '{filename}': {len(slides)} slides, theme={theme}")

            if not filename.endswith(".pptx"):
                filename += ".pptx"

            t = THEMES.get(theme.lower(), THEMES["blue"])
            prs = Presentation()
            prs.slide_width = Inches(13.33)
            prs.slide_height = Inches(7.5)

            for slide_data in slides:
                if not isinstance(slide_data, dict):
                    slide_data = {"type": "content", "title": str(slide_data)}
                slide_type = slide_data.get("type", "content").lower()
                builder = _BUILDERS.get(slide_type, _build_content_slide)
                builder(prs, slide_data, t)

            filepath = os.path.join(self.output_dir, filename)
            prs.save(filepath)
            url = f"http://{self.host}:{self.port}/outputs/{filename}"
            logger.info(f"Saved: {filepath}")
            return (
                f"\u2705 Successfully generated **{filename}**!\n\n"
                f"\U0001f4ca [Download your presentation here]({url})"
            )

        except Exception as exc:
            logger.error(f"PPTX generation error: {exc}", exc_info=True)
            return f"\u274c Failed to generate presentation: {exc}"

    def get_tools(self) -> dict[str, Any]:
        return {"generate_pptx": self}
