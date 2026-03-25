import asyncio
import json
import os
from pptx import Presentation
from typing import Any
from pydantic import BaseModel

class PPTXRequest(BaseModel):
    title: str
    slides: list[dict[str, str]]  # Each dict has "title" and "content"

class ComplianceToolset:
    """Toolset for generating professional PowerPoint presentations."""

    def __init__(self):
        self.output_dir = "outputs"
        os.makedirs(self.output_dir, exist_ok=True)

    async def create_presentation(self, presentation_title: str, slides_json: str) -> str:
        """
        Creates a .pptx file based on a title and a list of slide contents.
        
        Args:
            presentation_title: The main title of the PowerPoint.
            slides_json: A JSON string list of slides, e.g., '[{"title": "Intro", "content": "Hello"}]'
            
        Returns:
            str: A confirmation message with the filename.
        """
        try:
            slides_data = json.loads(slides_json)
            prs = Presentation()
            
            # Title Slide
            title_slide_layout = prs.slide_layouts[0]
            slide = prs.slides.add_slide(title_slide_layout)
            slide.shapes.title.text = presentation_title

            # Content Slides
            bullet_layout = prs.slide_layouts[1]
            for s in slides_data:
                slide = prs.slides.add_slide(bullet_layout)
                slide.shapes.title.text = s.get("title", "Untitled Slide")
                slide.placeholders[1].text = s.get("content", "")

            filename = f"{presentation_title.replace(' ', '_')}.pptx"
            filepath = os.path.join(self.output_dir, filename)
            prs.save(filepath)
            
            return f"Successfully created presentation: {filename} with {len(slides_data)} slides."
            
        except Exception as e:
            return f"Failed to create PPTX: {str(e)}"

    def get_tools(self) -> dict[str, Any]:
        """Registers the tools for OpenAI function calling."""
        return {
            'create_presentation': self.create_presentation
        }
