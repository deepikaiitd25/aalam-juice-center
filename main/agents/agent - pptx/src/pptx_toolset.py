import json
import logging
import os
import uuid
from typing import Any, Union

from pydantic import BaseModel
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

logger = logging.getLogger(__name__)

# ---------------------------------------------------------------------------
# Response model
# ---------------------------------------------------------------------------


class PptxGenerationResponse(BaseModel):
    status: str
    file_url: str | None = None
    error_message: str | None = None


# ---------------------------------------------------------------------------
# Theme definitions
# ---------------------------------------------------------------------------

THEMES: dict[str, dict] = {
    "blue": {
        "accent":  RGBColor(0x1F, 0x49, 0x7D),
        "light":   RGBColor(0xDE, 0xEB, 0xF7),
        "white":   RGBColor(0xFF, 0xFF, 0xFF),
        "dark":    RGBColor(0x26, 0x26, 0x26),
    },
    "green": {
        "accent":  RGBColor(0x17, 0x5E, 0x35),
        "light":   RGBColor(0xD5, 0xEB, 0xDA),
        "white":   RGBColor(0xFF, 0xFF, 0xFF),
        "dark":    RGBColor(0x22, 0x22, 0x22),
    },
    "dark": {
        "accent":  RGBColor(0x16, 0x21, 0x3E),
        "light":   RGBColor(0x53, 0x3A, 0x7A),
        "white":   RGBColor(0xF0, 0xF0, 0xF0),
        "dark":    RGBColor(0xEE, 0xEE, 0xEE),
    },
    "red": {
        "accent":  RGBColor(0x9B, 0x1C, 0x1C),
        "light":   RGBColor(0xF9, 0xDE, 0xDE),
        "white":   RGBColor(0xFF, 0xFF, 0xFF),
        "dark":    RGBColor(0x22, 0x22, 0x22),
    },
    "purple": {
        "accent":  RGBColor(0x4B, 0x0E, 0x82),
        "light":   RGBColor(0xE8, 0xD8, 0xF5),
        "white":   RGBColor(0xFF, 0xFF, 0xFF),
        "dark":    RGBColor(0x22, 0x22, 0x22),
    },
}


# ---------------------------------------------------------------------------
# Slide builder helpers
# ---------------------------------------------------------------------------

def _set_bg(slide, color: RGBColor):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def _rect(slide, left, top, width, height, color: RGBColor):
    shp = slide.shapes.add_shape(1, left, top, width, height)
    shp.fill.solid()
    shp.fill.fore_color.rgb = color
    shp.line.fill.background()
    return shp


def _textbox(slide, left, top, width, height, text, size, bold, color: RGBColor,
             align=PP_ALIGN.LEFT, wrap=True):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = wrap
    tf.text = text
    for para in tf.paragraphs:
        para.alignment = align
        for run in para.runs:
            run.font.size = Pt(size)
            run.font.bold = bold
            run.font.color.rgb = color
    return tb


def _build_title_slide(prs, data, t):
    W, H = prs.slide_width, prs.slide_height
    sl = prs.slides.add_slide(prs.slide_layouts[6])
    _set_bg(sl, t["accent"])
    _rect(sl, 0, 0, Inches(0.3), H, t["light"])
    _rect(sl, 0, H - Inches(1.6), W, Inches(1.6), t["white"])
    _textbox(sl, Inches(0.8), Inches(1.4), W - Inches(1.4), Inches(2.2),
             data.get("title", "Presentation"), 40, True, t["white"])
    sub = data.get("subtitle", "")
    if sub:
        _textbox(sl, Inches(0.8), Inches(3.8), W - Inches(1.4), Inches(1.0),
                 sub, 22, False, t["light"])


def _build_content_slide(prs, data, t):
    W, H = prs.slide_width, prs.slide_height
    sl = prs.slides.add_slide(prs.slide_layouts[6])
    _set_bg(sl, t["white"])
    hdr_h = Inches(1.3)
    _rect(sl, 0, 0, W, hdr_h, t["accent"])
    _textbox(sl, Inches(0.5), Inches(0.15), W - Inches(1.0), hdr_h,
             data.get("title", "Slide"), 28, True, t["white"])
    bullets = data.get("bullets", [])
    if bullets:
        tb = sl.shapes.add_textbox(Inches(0.6), hdr_h + Inches(0.3),
                                   W - Inches(1.2), H - hdr_h - Inches(0.8))
        tf = tb.text_frame
        tf.word_wrap = True
        for i, b in enumerate(bullets):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            p.text = f"•  {b}"
            p.space_before = Pt(7)
            for run in p.runs:
                run.font.size = Pt(18)
                run.font.color.rgb = t["dark"]
    notes = data.get("notes", "")
    if notes:
        sl.notes_slide.notes_text_frame.text = notes
    _rect(sl, 0, H - Inches(0.08), W, Inches(0.08), t["accent"])


def _build_two_column_slide(prs, data, t):
    W, H = prs.slide_width, prs.slide_height
    sl = prs.slides.add_slide(prs.slide_layouts[6])
    _set_bg(sl, t["white"])
    hdr_h = Inches(1.3)
    _rect(sl, 0, 0, W, hdr_h, t["accent"])
    _textbox(sl, Inches(0.5), Inches(0.15), W - Inches(1.0), hdr_h,
             data.get("title", "Comparison"), 28, True, t["white"])
    mid = W // 2
    col_top = hdr_h + Inches(0.25)
    col_h = H - col_top - Inches(0.5)
    col_w = mid - Inches(0.8)
    _rect(sl, mid - Inches(0.02), col_top, Inches(0.04), col_h, t["light"])

    def _col(left, col_title, col_bullets):
        _textbox(sl, left, col_top, col_w, Inches(0.5),
                 col_title, 19, True, t["accent"])
        if col_bullets:
            tb = sl.shapes.add_textbox(left, col_top + Inches(0.6),
                                       col_w, col_h - Inches(0.6))
            tf = tb.text_frame
            tf.word_wrap = True
            for i, b in enumerate(col_bullets):
                p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
                p.text = f"•  {b}"
                for run in p.runs:
                    run.font.size = Pt(17)
                    run.font.color.rgb = t["dark"]

    _col(Inches(0.5), data.get("left_title", "Option A"),
         data.get("left_bullets", []))
    _col(mid + Inches(0.3), data.get("right_title",
         "Option B"), data.get("right_bullets", []))
    _rect(sl, 0, H - Inches(0.08), W, Inches(0.08), t["accent"])


def _build_closing_slide(prs, data, t):
    W, H = prs.slide_width, prs.slide_height
    sl = prs.slides.add_slide(prs.slide_layouts[6])
    _set_bg(sl, t["accent"])
    _rect(sl, 0, 0, Inches(0.4), H, t["light"])
    _textbox(sl, Inches(1.2), Inches(1.8), W - Inches(2.0), Inches(2.4),
             data.get("title", "Thank You"), 48, True, t["white"],
             align=PP_ALIGN.CENTER)
    sub = data.get("subtitle", "")
    if sub:
        _textbox(sl, Inches(1.2), Inches(4.4), W - Inches(2.0), Inches(1.0),
                 sub, 22, False, t["light"], align=PP_ALIGN.CENTER)


# ---------------------------------------------------------------------------
# Toolset class
# ---------------------------------------------------------------------------

# ---------------------------------------------------------------------------
# Toolset class
# ---------------------------------------------------------------------------

class PptxToolset:
    """Toolset for generating PowerPoint presentations."""

    def __init__(self, host: str, port: int):
        self.host = host
        self.port = port
        self.output_dir = "outputs"
        os.makedirs(self.output_dir, exist_ok=True)
        logger.info(
            f"Initialized PptxToolset. Saving files to ./{self.output_dir}")

    # THE FIX: Make it async and return a clean text string
    async def generate_pptx(
        self,
        filename: str,
        slides: list,
        theme: str = "blue",
    ) -> str:
        """
        Generates a PowerPoint (.pptx) presentation from structured slide data.

        Args:
            filename: Output filename (without extension, e.g. "startup_pitch")
            slides: List of slide objects describing the presentation.
            theme: Color theme — one of: blue, green, dark, red, purple (default: blue)
        """
        try:
            if isinstance(slides, str):
                logger.info("slides received as string — parsing JSON...")
                slides = json.loads(slides)

            logger.info(
                f"Generating PPTX: {filename} with {len(slides)} slides, theme={theme}")

            if not filename.endswith(".pptx"):
                filename += ".pptx"

            t = THEMES.get(theme.lower(), THEMES["blue"])

            prs = Presentation()
            prs.slide_width = Inches(13.33)
            prs.slide_height = Inches(7.5)

            for slide_data in slides:
                # Fallback to empty dict if LLM hallucinates
                if not isinstance(slide_data, dict):
                    slide_data = {"type": "content", "title": str(slide_data)}

                slide_type = slide_data.get("type", "content").lower()

                if slide_type == "title":
                    _build_title_slide(prs, slide_data, t)
                elif slide_type == "two_column":
                    _build_two_column_slide(prs, slide_data, t)
                elif slide_type in ("closing", "end", "thank_you"):
                    _build_closing_slide(prs, slide_data, t)
                else:
                    _build_content_slide(prs, slide_data, t)

            filepath = os.path.join(self.output_dir, filename)
            prs.save(filepath)

            download_url = f"http://{self.host}:{self.port}/outputs/{filename}"
            logger.info(f"Saved PPTX to {filepath}")

            # THE FIX: Return standard Markdown for the Nasiko UI
            return f"✅ Successfully generated **{filename}**!\n\n📊 [Download your presentation here]({download_url})"

        except Exception as e:
            logger.error(f"Error generating PPTX: {e}")
            return f"❌ Failed to generate presentation: {str(e)}"

    def get_tools(self) -> dict[str, Any]:
        return {
            "generate_pptx": self,
        }
